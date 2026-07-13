import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

prs = Presentation('vllm-ascend-KV cache管理优化方案.pptx')

MAGENTA = RGBColor(198, 0, 84)
WHITE = RGBColor(255, 255, 255)
TEXT_PRIMARY = RGBColor(30, 30, 30)

# ====== IMPROVEMENT 1: Slide 23 (index 22) - Root cause fix ======
slide23 = prs.slides[22]
for shape in slide23.shapes:
    if shape.has_text_frame:
        tf = shape.text_frame
        ft = tf.text

        # Fix card title: "3. Sparse 多 Tensor" → "3. Spec 万能类: 类型擦除"
        if ft.startswith("3. Sparse"):
            for p in tf.paragraphs:
                if p.text.startswith("3. Sparse"):
                    for run in p.runs:
                        run.text = run.text.replace("3. Sparse 多 Tensor", "3. Spec 万能类: 类型擦除")
                        run.font.color.rgb = MAGENTA
                    break

        # Fix card body
        if "aclnnSparseFlashAttention" in ft or "3~4独立tensor输入" in ft:
            for p in tf.paragraphs:
                if "aclnnSparseFlashAttention" in p.text or "3~4独立tensor" in p.text or "host侧预拆分" in p.text:
                    for run in p.runs:
                        run.text = "5种MLA变体→1个AscendMLAAttentionSpec。isinstance()不编码语义，下游用if-else恢复已丢失的类型信息。"
                    break

        # Fix callout title
        if "两个根本源头" in ft:
            for p in tf.paragraphs:
                if "两个根本源头" in p.text:
                    for run in p.runs:
                        if "两个" in run.text:
                            run.text = run.text.replace("两个", "三个")
                    break

        # Fix callout body
        if "硬件约束: NPU算子不支持stride" in ft and "接口缺陷:" in ft:
            for p in tf.paragraphs:
                if p.text and "接口缺陷" in p.text:
                    for run in p.runs:
                        if "接口缺陷" in run.text:
                            run.text = (
                                "硬件约束: NPU算子不支持stride访问 → host侧必须预拆分tensor → 分配和重塑逻辑复杂化 (根因1,2,4,5)。"
                                "Spec 设计: 类型不编码语义 → isinstance()无区分力 → 下游被迫用if-else恢复丢失的类型信息 (根因3)。"
                                "接口缺陷: get_kv_cache_shape() 单返回值无法表达多tensor布局 → 信息泄漏到 model_runner (根因6)。"
                            )
                    break

    # Fix Layout refactor card
    if shape.has_text_frame:
        tf = shape.text_frame
        if "KVCacheLayout 抽象基类" in tf.text and "StandardKV" in tf.text:
            for p in tf.paragraphs:
                if "6 个子类" in p.text and "Spec" not in p.text:
                    for run in p.runs:
                        if "6 个子类" in run.text:
                            run.text = run.text.replace(
                                "6 个子类: StandardKV, MLANopeRope, SparseMLA,",
                                "Spec 4~5 子类 + Layout 6 个子类:\n  Spec: AscendStandardMLA/SparseMLA/SparseC8MLA/DSV4MLA\n  Layout: StandardKV, MLANopeRope, SparseMLA,"
                            )
                    break

print("Improvement 1 done: Slide 23 root cause updated.")

# ====== IMPROVEMENT 2: Slide 15 (index 14) - Rebuild table ======
slide15 = prs.slides[14]

for shape in slide15.shapes:
    if shape.has_table:
        tbl = shape.table
        header_cell = tbl.cell(0, 0).text
        if "模型类型" in header_cell or "模型" in header_cell:
            # Save position
            tbl_left = shape.left
            tbl_top = shape.top

            # Remove old table
            sp = shape._element
            sp.getparent().remove(sp)

            # New data with upstream columns
            new_data = [
                ["模型类型", "上游 Spec (基线)", "Ascend Spec",
                 "上游 Backend (基线)", "Ascend Backend",
                 "上游 Reshape (基线)", "Ascend Reshape", "差距 Δ"],
                ["GQA\n(Qwen3)",
                 "FullAttentionSpec\n7字段,零NPU",
                 "复用上游",
                 "FlashAttention\n→ (2,N,B,H,D)",
                 "AscendAttention\n→ (2,N,B,H,D)",
                 "view(N,B,H,D)\n1行,无分支",
                 ".view()\n同上游",
                 "★ K/V物理拆分"],
                ["标准MLA\n(DS V3.1)",
                 "MLAAttentionSpec\nhead_size=576",
                 "AscendMLAAttention\n+scale_dim/dtype",
                 "FlashMLA\n→ (N,B,576)",
                 "AscendMLA\n→ (N,B,1,576)",
                 "view(N,B,576)\n1行,无分支",
                 "view+查layer\nhead_size丢失512+64",
                 "★★ head_size=576\n丢失拆分信息"],
                ["SparseMLA\n(DS V3.2)",
                 "MLAAttentionSpec\n+cache_dtype_str",
                 "AscendMLAAttention\n+sparse_head_dim/C8",
                 "FlashMLASparse\n→ (N,B,576)",
                 "AscendSFA\n→ (N,B,1,576)",
                 "view(N,B,576)\n1行,无分支",
                 "unpack+view\nC8 x A5/A3多分支",
                 "★★★ indexer+C8\nA5/A3硬件分支"],
                ["CompressMLA\n(DS V4)",
                 "MLAAttentionSpec\n+compress_ratio",
                 "AscendMLAAttention\n+compress_ratio/scale",
                 "FlashMLASparse\n→ (N,B,656)",
                 "AscendDSA\n→ (N,B,1,head)",
                 "view(N,B,656)\n1行,无分支",
                 "as_strided overlay\n~100行,3 views",
                 "★★★★ 时间压缩+fp8\nas_strided承载语义"],
                ["Hybrid\n(Qwen3.5)",
                 "FullAttn+MambaSpec\n2种原生类型",
                 "复用上游\n+page_size_padded",
                 "FlashAttn+GDN\n→ 两种shape",
                 "AscendAttn+GDN\n→ (2,N,B,4,256)",
                 "view+as_strided_\n零拷贝重排,~5行",
                 "slice+strip+padding\n~100行hybrid专用",
                 "★★★★★ attn+mamba\nlayout互斥"],
            ]

            new_rows, new_cols = len(new_data), len(new_data[0])

            # Create new table
            new_shape = slide15.shapes.add_table(
                new_rows, new_cols, tbl_left, tbl_top,
                Inches(12.7), Inches(3.4))
            new_tbl = new_shape.table

            for r in range(new_rows):
                for c in range(new_cols):
                    cell = new_tbl.cell(r, c)
                    p = cell.text_frame.paragraphs[0]
                    p.text = new_data[r][c]
                    p.font.size = Pt(6.5)
                    p.font.color.rgb = TEXT_PRIMARY
                    p.font.name = "Consolas" if r > 0 else "Microsoft YaHei"
                    if r == 0:
                        p.font.bold = True
                        p.font.color.rgb = WHITE
                        p.font.size = Pt(7)
                    # Cell fill
                    tcPr = cell._tc.get_or_add_tcPr()
                    fill = cell._tc.makeelement(qn("a:solidFill"), {})
                    clr_val = "C6000B" if r == 0 else "F5F5F7"
                    clr = cell._tc.makeelement(qn("a:srgbClr"), {"val": clr_val})
                    fill.append(clr)
                    tcPr.insert(0, fill)

            print("Improvement 2 done: Slide 15 table rebuilt with upstream columns.")
            break

prs.save("vllm-ascend-KV cache管理优化方案.pptx")
print(f"Saved successfully. Total slides: {len(prs.slides)}")
