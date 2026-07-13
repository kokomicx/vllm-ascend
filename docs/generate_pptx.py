#!/usr/bin/env python3
"""PPTX generator — 总体→局部 叙事结构: 端到端管线 → 逐阶段深入 → 根因与对齐."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ═══════════════════════════════════════════
# COMPANY BRAND COLOR SCHEME
# ═══════════════════════════════════════════
PRIMARY_RED   = RGBColor(198, 0, 11)
DEEP_RED      = RGBColor(200, 16, 46)
RED_TINT1     = RGBColor(220, 80, 85)
RED_TINT2     = RGBColor(235, 145, 148)
RED_TINT3     = RGBColor(248, 210, 212)
GREEN         = RGBColor(88, 178, 48)
YELLOW        = RGBColor(252, 200, 0)
ORANGE        = RGBColor(237, 108, 0)
CYAN          = RGBColor(43, 181, 197)
MAGENTA       = RGBColor(198, 0, 84)
BLACK         = RGBColor(0, 0, 0)
DARK_GRAY     = RGBColor(89, 87, 87)
MID_GRAY      = RGBColor(137, 137, 137)
LIGHT_GRAY    = RGBColor(200, 200, 200)
OFF_WHITE     = RGBColor(245, 245, 247)
WHITE         = RGBColor(255, 255, 255)

BG_DEEP      = WHITE
BG_PRIMARY   = WHITE
BG_CARD      = OFF_WHITE
BORDER       = RGBColor(218, 218, 222)
TEXT_PRIMARY = RGBColor(30, 30, 30)
TEXT_SEC     = DARK_GRAY
TEXT_MUTED   = MID_GRAY
BLUE         = PRIMARY_RED
AMBER        = ORANGE
RED          = DEEP_RED
PURPLE       = MAGENTA

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ═══════════════════════════════════════════
# HELPERS
# ═══════════════════════════════════════════

def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def tb(slide, l, t, w, h, text, sz=Pt(14), color=TEXT_PRIMARY, bold=False,
       align=PP_ALIGN.LEFT, font='Microsoft YaHei', ls=1.2):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = sz
    p.font.color.rgb = color; p.font.bold = bold; p.font.name = font
    p.alignment = align; p.line_spacing = Pt(sz.pt * ls)
    return tf

def multi_tb(slide, l, t, w, h, lines, sz=Pt(9), color=TEXT_SEC, lh=16):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = True; tf.paragraphs[0].text = ""
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = sz; p.font.color.rgb = color
        p.font.name = 'Microsoft YaHei'
        p.space_before = Pt(0); p.space_after = Pt(0); p.line_spacing = Pt(lh)
    return tf

def code_box(slide, l, t, w, lines, sz=Pt(8)):
    CODE_BG = RGBColor(35, 35, 40); CODE_TEXT = RGBColor(210, 210, 215)
    lh_pt = sz.pt * 1.5; total_h = len(lines) * Pt(int(lh_pt)) + Pt(12)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, total_h)
    shape.fill.solid(); shape.fill.fore_color.rgb = CODE_BG
    shape.line.color.rgb = RGBColor(65, 65, 70); shape.line.width = Pt(0.5)
    shape.adjustments[0] = 0.03
    tf = shape.text_frame; tf.word_wrap = True; tf.paragraphs[0].text = ""
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = sz; p.font.color.rgb = CODE_TEXT
        p.font.name = 'Consolas'
        p.space_before = Pt(0); p.space_after = Pt(0); p.line_spacing = Pt(lh_pt)
    return shape

def arch_box(slide, l, t, w, lines, sz=Pt(8)):
    """Architecture diagram box (monospace, card bg)."""
    lh_pt = sz.pt * 1.55
    total_h = len(lines) * Pt(int(lh_pt)) + Pt(14)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, total_h)
    shape.fill.solid(); shape.fill.fore_color.rgb = BG_CARD
    shape.line.color.rgb = BORDER; shape.line.width = Pt(1)
    shape.adjustments[0] = 0.03
    tf = shape.text_frame; tf.word_wrap = False; tf.paragraphs[0].text = ""
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = sz; p.font.color.rgb = TEXT_PRIMARY
        p.font.name = 'Consolas'
        p.space_before = Pt(0); p.space_after = Pt(0); p.line_spacing = Pt(lh_pt)
    return shape

def card(slide, l, t, w, h, title, lines, tc=BLUE, accent=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    shape.fill.solid(); shape.fill.fore_color.rgb = BG_CARD
    shape.line.color.rgb = BORDER; shape.line.width = Pt(1); shape.adjustments[0] = 0.04
    if accent:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, Pt(3), h)
        bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background()
    tb(slide, l+Pt(12), t+Pt(8), w-Pt(24), Pt(20), title, sz=Pt(12), color=tc, bold=True)
    multi_tb(slide, l+Pt(12), t+Pt(30), w-Pt(24), h-Pt(38), lines, sz=Pt(8.5), lh=15)

def callout(slide, l, t, w, label, body, accent=AMBER):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, Pt(48))
    r = max(0, min(255, int(accent[0]*0.08 + BG_PRIMARY[0]*0.92)))
    g = max(0, min(255, int(accent[1]*0.08 + BG_PRIMARY[1]*0.92)))
    b = max(0, min(255, int(accent[2]*0.08 + BG_PRIMARY[2]*0.92)))
    shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(r,g,b)
    shape.line.color.rgb = RGBColor(min(accent[0]+60,255), min(accent[1]+60,255), min(accent[2]+60,255))
    shape.line.width = Pt(1); shape.adjustments[0] = 0.04
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, Pt(4), Pt(48))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background()
    tb(slide, l+Pt(14), t+Pt(6), w-Pt(28), Pt(16), f"> {label}", sz=Pt(9), color=accent, bold=True)
    tb(slide, l+Pt(14), t+Pt(22), w-Pt(28), Pt(26), body, sz=Pt(8.5), color=TEXT_SEC)

def tag(slide, l, t, text, color):
    tw = Pt(len(text)*9+16)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, tw, Pt(18))
    r = max(0, min(255, int(color[0]*0.13+BG_PRIMARY[0]*0.87)))
    g = max(0, min(255, int(color[1]*0.13+BG_PRIMARY[1]*0.87)))
    b = max(0, min(255, int(color[2]*0.13+BG_PRIMARY[2]*0.87)))
    shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(r,g,b); shape.line.fill.background()
    shape.adjustments[0] = 0.5
    p = shape.text_frame.paragraphs[0]; p.text = text; p.font.size = Pt(8)
    p.font.color.rgb = color; p.font.bold = True
    p.font.name = 'Microsoft YaHei'; p.alignment = PP_ALIGN.CENTER

def stat(slide, l, t, num, label, nc=BLUE):
    tb(slide, l, t, Inches(2.2), Pt(36), num, sz=Pt(34), color=nc, bold=True, font='Consolas', align=PP_ALIGN.CENTER)
    tb(slide, l, t+Pt(36), Inches(2.2), Pt(30), label, sz=Pt(9), color=TEXT_SEC, align=PP_ALIGN.CENTER)

def slide_num(slide, n, total=25):
    tb(slide, Inches(11.6), Inches(0.25), Inches(1.4), Pt(16), f"{n} / {total}",
       sz=Pt(9), color=TEXT_MUTED, align=PP_ALIGN.RIGHT, font='Consolas')

def s_marker(slide, text):
    tb(slide, Inches(0.6), Inches(0.25), Inches(3), Pt(14), text, sz=Pt(9), color=BLUE, bold=True, font='Consolas')

def divider_slide(n, total, label, title, sub):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, BG_PRIMARY); slide_num(s, n, total)
    s_marker(s, label)
    tb(s, Inches(0.6), Inches(2.0), Inches(12), Inches(0.8), title, sz=Pt(32), color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    tb(s, Inches(0.6), Inches(3.0), Inches(12), Inches(0.4), sub, sz=Pt(13), color=TEXT_SEC, align=PP_ALIGN.CENTER)
    return s

def pipeline_flow(slide, l, t, w, stages, sz=Pt(7.5)):
    """Draw a horizontal pipeline: list of (label, color) tuples."""
    n = len(stages); box_w = (w - Pt((n-1)*8)) // n
    for i, (label, clr) in enumerate(stages):
        x = l + i * (box_w + Pt(8))
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, t, box_w, Pt(24))
        shape.fill.solid()
        r2 = max(0, min(255, int(clr[0]*0.1))); g2 = max(0, min(255, int(clr[1]*0.1))); b2 = max(0, min(255, int(clr[2]*0.1)))
        shape.fill.fore_color.rgb = RGBColor(r2,g2,b2)
        shape.line.color.rgb = clr; shape.line.width = Pt(1); shape.adjustments[0] = 0.3
        p = shape.text_frame.paragraphs[0]; p.text = label; p.font.size = sz
        p.font.color.rgb = clr; p.font.bold = True; p.font.name = 'Consolas'
        p.alignment = PP_ALIGN.CENTER
        if i < n-1:
            tb(slide, x+box_w-Pt(2), t-Pt(2), Pt(12), Pt(28), "→", sz=Pt(10), color=TEXT_MUTED, align=PP_ALIGN.CENTER)

def make_table(slide, l, t, w, h, data, header_color='C6000B', row_color='F5F5F7'):
    """Create a styled table from 2D list of strings."""
    rows, cols = len(data), len(data[0])
    tbl = slide.shapes.add_table(rows, cols, l, t, w, h).table
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c); p = cell.text_frame.paragraphs[0]
            p.text = data[r][c]; p.font.size = Pt(7.5); p.font.color.rgb = TEXT_PRIMARY
            p.font.name = 'Microsoft YaHei' if r == 0 else 'Consolas'
            if r == 0: p.font.bold = True; p.font.color.rgb = WHITE; p.font.size = Pt(8)
            tcPr = cell._tc.get_or_add_tcPr()
            fill = cell._tc.makeelement(qn('a:solidFill'), {})
            clr = cell._tc.makeelement(qn('a:srgbClr'), {'val': header_color if r == 0 else row_color})
            fill.append(clr); tcPr.insert(0, fill)
    return tbl

TOTAL = 30

# ═══════════════════════════════════════════
# SLIDE 1 — COVER
# ═══════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_PRIMARY); slide_num(s, 1, TOTAL)

# Red accent bar at top
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Pt(5))
bar.fill.solid(); bar.fill.fore_color.rgb = PRIMARY_RED; bar.line.fill.background()

tag(s, Inches(4.0), Inches(1.5), "社区 vLLM", GREEN)
tag(s, Inches(5.2), Inches(1.5), "vLLM-Ascend", AMBER)
tag(s, Inches(6.7), Inches(1.5), "技术汇报", BLUE)

tb(s, Inches(1.5), Inches(2.2), Inches(10.3), Inches(1.2),
   "KV Cache 管理方式\n梳理与差异分析", sz=Pt(40), color=PRIMARY_RED, bold=True, align=PP_ALIGN.CENTER, ls=1.15)
tb(s, Inches(1.5), Inches(3.5), Inches(10.3), Inches(0.4),
   "从端到端管线 (Spec→Backend→Allocate→Reshape) 到逐阶段代码级对比", sz=Pt(14), color=TEXT_SEC, align=PP_ALIGN.CENTER)
tb(s, Inches(1.5), Inches(4.0), Inches(10.3), Inches(0.3),
   "覆盖: GQA(Qwen3 MoE) / MLA(DS V3.1/V3.2/V4) / SFA(GLM5.1) / Hybrid(Qwen3.5)",
   sz=Pt(11), color=TEXT_SEC, align=PP_ALIGN.CENTER)

stat(s, Inches(2.2), Inches(4.8), "~600 行", "Ascend model_runner\n分支代码", RED)
stat(s, Inches(4.8), Inches(4.8), "~60 行", "上游 model_runner\n代码量", GREEN)
stat(s, Inches(7.4), Inches(4.8), "15+", "if-else 分支点", BLUE)
stat(s, Inches(10.0), Inches(4.8), "5", "模型类型覆盖", PURPLE)

tb(s, Inches(1.5), Inches(6.5), Inches(10.3), Inches(0.3),
   "2026 年 7 月 · 内部技术评审", sz=Pt(10), color=TEXT_MUTED, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════
# SLIDE 2 — 汇报路线图 (取代传统 Agenda)
# ═══════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_PRIMARY); slide_num(s, 2, TOTAL)
tb(s, Inches(0.6), Inches(0.3), Inches(12), Inches(0.5), "汇报路线图：从总体管线到逐阶段代码级对比", sz=Pt(26), color=BLUE, bold=True)

roadmap = [
    ("Part 1", "总体架构", "端到端管线全景：Spec → Backend → Allocate → Reshape\n社区 KV Cache 管理方法速览", BLUE, "p1"),
    ("Part 2", "Spec 层", "上游「类型即语义」vs Ascend「万能类」\n以 DS V3.2 为例的 Spec 创建代码对比", CYAN, "p2"),
    ("Part 3", "Backend 层", "4 类 Ascend Backend vs 上游 Backend\nget_kv_cache_shape() 的 1/N 信息问题", GREEN, "p3"),
    ("Part 4", "Allocate & Reshape", "上游 30 行 vs Ascend 185 行\nDS V3.2 / DS V4 代码拆解", AMBER, "p4"),
    ("Part 5", "模型类型差异", "5 种模型类型完整链路对比表\n标准 MLA → Sparse MLA → Compress MLA", PURPLE, "p5"),
    ("Part 6", "精度 + 根因 + 对齐", "bf16 vs fp8 正交维度 · 六大根因\nLayout 重构方案 · 三步路线图", RED, "p6"),
]
for i, (part, title, desc, clr, _) in enumerate(roadmap):
    y = Inches(1.1) + Inches(i * 1.0)
    # Part number badge
    badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), y, Inches(1.3), Pt(24))
    r2 = max(0, min(255, int(clr[0]*0.1))); g2 = max(0, min(255, int(clr[1]*0.1))); b2 = max(0, min(255, int(clr[2]*0.1)))
    badge.fill.solid(); badge.fill.fore_color.rgb = RGBColor(r2,g2,b2); badge.line.fill.background()
    badge.adjustments[0] = 0.5
    p = badge.text_frame.paragraphs[0]; p.text = part; p.font.size = Pt(8); p.font.color.rgb = clr
    p.font.bold = True; p.font.name = 'Microsoft YaHei'; p.alignment = PP_ALIGN.CENTER
    # Title and description
    tb(s, Inches(2.1), y+Pt(0), Inches(3), Pt(22), title, sz=Pt(13), color=clr, bold=True)
    tb(s, Inches(5.3), y+Pt(0), Inches(7.5), Pt(38), desc, sz=Pt(8.5), color=TEXT_SEC)

callout(s, Inches(0.5), Inches(7.1), Inches(12.3), "叙事逻辑",
        "总体→局部: 先建立「Spec→Backend→Allocate→Reshape」端到端管线共识 → 再逐阶段深入 + 配具体模型代码 → 最后归纳根因与对齐方案。每个阶段都回答同一个问题: 上游怎么做? Ascend 怎么做? 差距在哪?", BLUE)

# ═══════════════════════════════════════════
# PART 1: 总体架构 — 端到端管线全景
# ═══════════════════════════════════════════
divider_slide(3, TOTAL, "PART 1 · 总体架构", "端到端管线全景：Spec → Backend → Allocate → Reshape",
              "先建立全局共识——上游和Ascend各自怎么走完这四步——再逐阶段深入")

# SLIDE 4 — 上游社区 vLLM: 端到端管线 (5 阶段)
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_PRIMARY); slide_num(s, 4, TOTAL)
tb(s, Inches(0.6), Inches(0.25), Inches(12), Inches(0.5), "社区 vLLM 端到端管线：四阶段，Backend 说了算", sz=Pt(24), color=GREEN, bold=True)
tb(s, Inches(0.6), Inches(0.7), Inches(12), Inches(0.25),
   "model_runner 的角色是「机械执行」——它不做模型类型的判断，只按 Backend 返回的结果做 view()/permute()。所有模型类型都走同一条路径。",
   sz=Pt(9), color=TEXT_SEC)

# Pipeline flow
pipeline_flow(s, Inches(0.5), Inches(1.15), Inches(12.3), [
    ("① Spec 创建", BLUE),
    ("② Backend 返回 Shape", CYAN),
    ("③ 全局规划: configs", GREEN),
    ("④ Allocate: zeros(int8)", AMBER),
    ("⑤ Reshape: view+permute", PURPLE),
], sz=Pt(8))

# Stage detail cards
stages_upstream = [
    ("① Spec 创建: layer → KVCacheSpec", GREEN, [
        "Attention.get_kv_cache_spec() 返回 frozen dataclass",
        "每层声明: block_size, num_kv_heads, head_size, dtype, page_size_bytes",
        "所有模型类型共用 KVCacheSpec 子类体系",
        "→ 对 model_runner 来说: 拿到一个 Spec 对象，不用问是什么模型",
    ]),
    ("② Backend 返回 Shape: 完整信息", CYAN, [
        "Backend.get_kv_cache_shape() 返回完整的 shape tuple",
        "GQA: (2, N, B, H, D) — K/V 交织在 dim=0",
        "MLA: (N, B, 576) — kv_lora(512)+k_rope(64) 平面存储",
        "→ 一个 shape 回答全部存储问题，1 tensor / 层",
    ]),
    ("③ 全局规划: get_kv_cache_configs()", BLUE, [
        "分组: get_kv_cache_groups() — 决策级联 (Uniform→Hybrid→DSV4)",
        "算块: available_memory / page_size_bytes → num_blocks",
        "出配置: KVCacheConfig { num_blocks, kv_cache_tensors[], groups[] }",
        "→ KVCacheTensor.shared_by 实现多层共享物理 buffer",
    ]),
    ("④⑤ Allocate + Reshape: ~60 行，0 分支", AMBER, [
        "_allocate_kv_cache(): torch.zeros(size, dtype=int8) — 30 行",
        "_reshape_kv_cache(): raw.view(dtype).view(shape).permute() — 30 行",
        "Mamba/MLA 用 as_strided 做零拷贝视图",
        "→ model_runner 不关心内部数据怎么分——全由 Backend 封装",
    ]),
]
for i, (title, color, lines) in enumerate(stages_upstream):
    x = Inches(0.3) + Inches(i * 3.2)
    card(s, x, Inches(1.7), Inches(3.1), Inches(2.6), title, lines, tc=color, accent=color)

# Key takeaway
callout(s, Inches(0.3), Inches(4.5), Inches(12.4), "社区核心原则",
        "Backend 是 KV Cache shape 的唯一决定者。每一维的含义、tensor 的数量、K/V 的存储方式，都由 Backend 说了算。model_runner 只做机械执行——拿到 raw tensor → view(dtype) → view(shape) → permute() → 结束。", GREEN)

# Additional detail cards below
card(s, Inches(0.3), Inches(5.3), Inches(6.0), Inches(1.4),
     "内部数据组织: 单 tensor 内部分拆", [
         "GQA: tensor[0]=K, tensor[1]=V — kernel 通过 offset 区分",
         "MLA: tensor[:,:,:512]=kv_lora, tensor[:,:,512:]=k_rope",
         "MLA fp8: 单 tensor, 字节级自定义布局 (656B/token V3.2)",
         "→ kernel 内部解析，host 侧完全不需要知道内部结构",
     ], tc=GREEN, accent=GREEN)

card(s, Inches(6.6), Inches(5.3), Inches(6.1), Inches(1.4),
     "分组方案: get_kv_cache_groups() 决策级联", [
         "情况1: 所有层相同 → 1 group 全包含",
         "情况2: 同类型不同维度 → UniformTypeKVCacheSpecs",
         "情况3: DS V4 → 按 page_size 分桶 + layer_tuple 对齐",
         "情况4: 通用 Hybrid → min_num_layers 分组 + padding 补齐",
     ], tc=GREEN, accent=GREEN)

# ═══════════════════════════════════════════
# SLIDE 5 — Ascend 的同一管线：问题在哪里
# ═══════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_PRIMARY); slide_num(s, 5, TOTAL)
tb(s, Inches(0.6), Inches(0.25), Inches(12), Inches(0.5), "Ascend 的同一管线：为什么变成了 600 行 if-else", sz=Pt(24), color=RED, bold=True)
tb(s, Inches(0.6), Inches(0.7), Inches(12), Inches(0.25),
   "上游的四个阶段在 Ascend 侧全部出现了信息丢失——每一步都只能完成 1/N 的工作，剩余 N-1/N 泄漏到 model_runner。",
   sz=Pt(9), color=TEXT_SEC)

pipeline_flow(s, Inches(0.5), Inches(1.15), Inches(12.3), [
    ("① Spec 创建", RED),
    ("② Backend: 只返参考Shape", RED),
    ("③ 全局规划: 有 patch", RED),
    ("④ Allocate: 多tensor拆分", RED),
    ("⑤ Reshape: 15+ if-else", RED),
], sz=Pt(8))

stages_ascend = [
    ("① Spec: AscendMLAAttentionSpec「万能类」", RED, [
        "上游: 1 种模型 → 1 种 Spec 子类 (类型即语义)",
        "Ascend: 所有 MLA 变体挤进同一个 AscendMLAAttentionSpec",
        "新增 6 个 NPU 特有字段: scale_dim, scale_dtype,",
        "  sparse_head_dim, cache_sparse_c8, c8_k_cache_dtype...",
        "→ Spec 声明了信息，但缺少解释方法 (page_size_bytes 有分支)",
    ]),
    ("② Backend: 只返回「参考 Shape」— 1/N 信息", RED, [
        "上游: get_kv_cache_shape() → (2,N,B,H,D) 完整可用",
        "Ascend: 同名方法 → (N,B,1,576) 参考值",
        "Sparse MLA 实际需要 3~4 tensor，Backend 只返回 1 个 shape",
        "→ 剩余 N-1 个 tensor 的 shape 需要 model_runner 自己推断",
    ]),
    ("③④⑤ Allocate+Reshape: 模型类型判断泄漏", RED, [
        "model_runner 被迫检查: use_sparse? use_compress? use_mla?",
        "  use_hybrid_blocks? C8? A5? A3? cache_only?",
        "Allocate: ~185 行, 判断需要拆几个 tensor + 各自多大",
        "Reshape: ~350 行, 判断每个 tensor 的 dtype + shape",
        "→ 15+ if-else, 每个新模型加 3~5 处分支",
    ]),
]
for i, (title, color, lines) in enumerate(stages_ascend):
    x = Inches(0.3) + Inches(i * 4.2)
    card(s, x, Inches(1.7), Inches(4.1), Inches(2.6), title, lines, tc=color, accent=color)

callout(s, Inches(0.3), Inches(4.5), Inches(12.4), "一句话总结核心差距",
        "上游 Backend 的 get_kv_cache_shape() 一个方法回答了全部存储问题；Ascend Backend 的同名方法只回答了 1/N 的信息 (一个 tensor 的 shape)，剩下 N-1/N 泄漏到了 model_runner 的 15+ 个 if-else 分支里。", RED)

# Bottom comparison cards
card(s, Inches(0.3), Inches(5.3), Inches(6.0), Inches(1.4),
     "根因 1: NPU 算子不支持 stride", [
         "CUDA kernel: SIMT LSU 可处理任意 stride 访问",
         "→ 单 tensor (2,N,B,H,D)，kernel 内部 offset 切 K/V",
         "Ascend NPU: DMA Engine 只接受 {base_addr, length}",
         "→ host 侧必须预拆分为独立连续 tensor → 分配+重塑复杂化",
     ], tc=RED, accent=RED)

card(s, Inches(6.6), Inches(5.3), Inches(6.1), Inches(1.4),
     "根因 2: Backend 接口语义不足", [
         "上游: tuple[int,...] 返回值 → 恰好够用 (1 tensor)",
         "Ascend: 需要表达 2~4 个 tensor 的布局",
         "→ 单返回值语义空间不够 → 只能返回「参考 shape」",
         "→ 拆分信息 (split ratio, dtype, count) 泄漏到外部",
     ], tc=RED, accent=RED)

# ═══════════════════════════════════════════
# SLIDE 6 — 社区 KV Cache 管理（一）：分层分组方案
# ═══════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_PRIMARY); slide_num(s, 6, 30)
tb(s, Inches(0.6), Inches(0.2), Inches(12), Inches(0.5), "社区 KV Cache 管理（一）：分层分组方案 (Grouping Schemes)", sz=Pt(22), color=GREEN, bold=True)
tb(s, Inches(0.6), Inches(0.65), Inches(12), Inches(0.25),
   "get_kv_cache_groups() 决策级联——从简到繁依次尝试，上游用一套算法覆盖所有模型类型（纯GQA / Hybrid / DS V4 / Mamba混合）。",
   sz=Pt(9), color=TEXT_SEC)

# Decision cascade as flow cards
card(s, Inches(0.3), Inches(1.1), Inches(3.0), Inches(2.1),
     "情况1: Uniform Spec", [
         "所有层 KVCacheSpec 完全相同",
         "→ 1 个 group，包含所有层",
         "适用: Llama, Qwen3 MoE (纯 GQA)",
         "page_size = num_layers ×",
         "  block_size × kv_hidden_size",
     ], tc=GREEN, accent=GREEN)

card(s, Inches(3.5), Inches(1.1), Inches(3.0), Inches(2.1),
     "情况2: Uniform Type", [
         "同 attention 类型，不同 hidden_size",
         "→ 1 个 UniformTypeKVCacheSpecs",
         "适用: 同类型但维度不同的层",
         "page_size = max(page_sizes)",
         "  通过 GCD 近似对齐",
     ], tc=BLUE, accent=BLUE)

card(s, Inches(6.7), Inches(1.1), Inches(3.1), Inches(2.1),
     "情况3: DeepSeek V4", [
         "SWA + MLA 三层结构",
         "→ group_and_unify + uniform groups",
         "按 page_size 分桶 + layer_tuple 对齐",
         "每 (tuple_idx, bucket) → KVCacheTensor",
         "num_blocks = mem / (tuple_bytes × tuples)",
     ], tc=PURPLE, accent=PURPLE)

card(s, Inches(10.0), Inches(1.1), Inches(3.0), Inches(2.1),
     "情况4: 通用 Hybrid", [
         "多 attention 类型混合",
         "→ min_num_layers 分组 + padding",
         "例: Gemma3 52sw+10full → 7 groups",
         "例: Llama4 3local+1full → 4 groups",
         "例: Qwen3.5 attn+mamba → 4 groups",
     ], tc=AMBER, accent=AMBER)

# Hybrid algorithm detail
code_box(s, Inches(0.3), Inches(3.45), Inches(6.2), [
    "# 通用 Hybrid 分组算法 (kv_cache_utils.py:1057-1176)",
    "# Step 1: 相同 spec 类型聚拢",
    "same_type_layers = defaultdict(list)",
    "for name, spec in kv_cache_spec.items():",
    "    same_type_layers[spec].append(name)",
    "# Step 2: 确定 group_size = min(每类层数)",
    "group_size = min(len(layers) for layers",
    "                  in same_type_layers.values())",
    "# Step 3: 每类层按 layers[i::num_groups] 交错分配",
    "#   (兼容 PP 跨 stage 对齐, 避免空 group)",
    "# Step 4: 不足 group_size 的用 padding 补齐",
    "#   → 创建 KVCacheGroupSpec",
], sz=Pt(7))

# Constraints table
group_data = [
    ["设计约束", "说明"],
    ["同 group 内相同 attention 类型", "保证 block 分配数量一致"],
    ["跨 group 相同 page_size", "内存池只有一个 page size"],
    ["每类 group 层数相等 (或加 padding)", "保证 block table 可重复"],
    ["每 token 每层物理内存相同", "当前强假设, Mamba 例外需 padding"],
    ["KV Sharing 层不参与分配", "直接复用目标层 KV cache (gemma-3n)"],
]
make_table(s, Inches(6.8), Inches(3.45), Inches(6.0), Inches(1.8), group_data, header_color='188C30')

# Key insight
callout(s, Inches(0.3), Inches(5.5), Inches(12.7), "分组层的核心设计思想",
        "将不同 attention 类型的层按「最小公倍数」分组，确保每个 group 有相同数量的层、所有 group 有相同的 page_size。这样上层 BlockPool 可以用统一的 block 大小管理所有 KV cache，避免内存碎片。Mamba 模型通过增大 attention 的 block_size 来对齐两种不同 state_size。", GREEN)

# ═══════════════════════════════════════════
# SLIDE 7 — 社区 KV Cache 管理（二）：空间分配与大小计算
# ═══════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_PRIMARY); slide_num(s, 7, 30)
tb(s, Inches(0.6), Inches(0.2), Inches(12), Inches(0.5), "社区 KV Cache 管理（二）：空间分配与大小计算", sz=Pt(22), color=GREEN, bold=True)
tb(s, Inches(0.6), Inches(0.65), Inches(12), Inches(0.25),
   "一切从 page_size_bytes 开始——一个 block 占多少字节。不同 Spec 子类独立实现，无 if-else 分支。",
   sz=Pt(9), color=TEXT_SEC)

# Left: page_size_bytes computation tree
arch_box(s, Inches(0.3), Inches(1.1), Inches(6.0), [
    "page_size_bytes 计算层级树 (每个子类独立 override)",
    "",
    "AttentionSpec.page_size_bytes (顶层包装)",
    "  = real_page_size_bytes",
    "  + FP8_PER_TOKEN_HEAD: 2×B×H×sizeof(f32)  ← K/V scale 预留",
    "  × page_size_padded 对齐",
    "",
    "  ├── FullAttentionSpec.real_page_size_bytes",
    "  │     = B × H × (head_size + head_size_v) × dtype_size",
    "  │     NVFP4: B × H × (head_size//2 + head_size//16)",
    "  │",
    "  ├── MLAAttentionSpec.real_page_size_bytes",
    "  │     fp8_ds_mla+V4: storage_B × 584  (448+128+8)",
    "  │     fp8_ds_mla+V3.2: B × 656       (512+128+16)",
    "  │     default: storage_B × H × head_size × dtype_size",
    "  │     storage_B = B // compress_ratio",
    "  │",
    "  └── MambaSpec.page_size_bytes",
    "        = sum(prod(shape) × dtype_size for each state)",
])

# Right: allocation flow + formula examples
card(s, Inches(6.6), Inches(1.1), Inches(6.2), Inches(1.7),
     "从 available_memory 到 num_blocks", [
         "情况 A (Uniform): num_blocks = available // page_size_bytes",
         "  每层一个 KVCacheTensor",
         "情况 B (DS V4): num_blocks = available // (tuple_bytes × num_tuples)",
         "  每个 (tuple_idx, bucket) 一个 KVCacheTensor",
         "情况 C (Hybrid): num_blocks = available // page_size // group_size",
         "  多个 group 的对应层 shared_by 合并",
         "",
         "最终产物: KVCacheConfig { num_blocks, kv_cache_tensors[], groups[] }",
         "Worker 侧: torch.zeros(tensor.size, dtype=int8)  ← 30行, 0分支!",
     ], tc=BLUE, accent=BLUE)

# Allocation formula examples
alloc_data = [
    ["模型 (B=16)", "Spec 类型", "page_size_bytes 公式", "实际值", "1 token bytes"],
    ["Qwen3 MoE bf16", "FullAttention", "16×8×(128+128)×2", "65,536 B", "256 B"],
    ["DS V3.1 MLA bf16", "MLAAttention", "16×1×576×2", "18,432 B", "1,152 B"],
    ["DS V3.2 fp8", "MLAAttention", "16×656", "10,496 B", "656 B"],
    ["DS V4 fp8 C4", "MLAAttention", "(16//4)×584", "2,336 B", "584/C B"],
    ["Qwen3.5 attn bf16", "FullAttention", "16×4×(256+256)×2", "65,536 B", "4,096 B"],
    ["Qwen3.5 mamba", "MambaSpec", "sum(shapes×dtype)", "取决于 state", "state_size"],
    ["NVFP4 GQA", "FullAttention", "16×8×(64+8)", "9,216 B", "72 B"],
]
make_table(s, Inches(0.3), Inches(3.05), Inches(12.4), Inches(2.15), alloc_data, header_color='188C30')

# per-request max memory
card(s, Inches(0.3), Inches(5.4), Inches(6.0), Inches(1.3),
     "per-request 最大内存: max_memory_usage_bytes()", [
         "FullAttn: cdiv(max_len, B) × page_size_bytes",
         "  // (dcp_world × pcp_world)  # 上下文并行",
         "SlidingWindow: min(sw-1+tokens, max_len) token 上限",
         "Mamba: 'all'→cdiv(max_len,B) / 'align'→2B / 'none'→1B",
         "UniformTypeKVCacheSpecs: max(内部各 spec 的计算值)",
     ], tc=BLUE, accent=BLUE)

card(s, Inches(6.6), Inches(5.4), Inches(6.2), Inches(1.3),
     "KVCacheTensor.shared_by: 多层共享物理 buffer", [
         "KVCacheConfig 的 kv_cache_tensors 列表中的每个元素:",
         "  size: 总字节数 = page_size_bytes × num_blocks",
         "  shared_by: [layer_name, ...] ← 这些层共享此 buffer",
         "",
         "例如 Qwen3.5 (4 groups × 8 层/group):",
         "  KVCacheTensor #0 shared_by = [layer.3.attn, layer.0.lin,",
         "                                 layer.8.lin, layer.16.lin]",
         "  ← 1 个 attn + 3 个 mamba 层共享同一个 int8 buffer!",
     ], tc=GREEN, accent=GREEN)

# ═══════════════════════════════════════════
# SLIDE 8 — 社区 KV Cache 管理（三）：内部数据组织与总 Shape
# ═══════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_PRIMARY); slide_num(s, 8, 30)
tb(s, Inches(0.6), Inches(0.2), Inches(12), Inches(0.5), "社区 KV Cache 管理（三）：内部数据组织与总 Shape", sz=Pt(22), color=GREEN, bold=True)
tb(s, Inches(0.6), Inches(0.65), Inches(12), Inches(0.25),
   "上游所有模型类型都产生 1 个 tensor/层（或 1 个 buffer + as_strided views）。不同类型数据如何组织，全部由 Backend 内部封装。",
   sz=Pt(9), color=TEXT_SEC)

# Main layout table
layout_full = [
    ["模型类型", "Tensor数", "总 Shape (逻辑)", "每 Token\n物理字节", "内部不同类型数据如何存储"],
    ["GQA\n(FlashAttn)", "1", "(2, N, B, H, D)", "H×(D+Dv)\n×dtype", "dim=0: [0]=K所有block, [1]=V所有block\nCUDA kernel 通过 base+offset 区分 K/V 区域\n→ host 侧完全不需要知道内部结构"],
    ["GQA\n(FlashInfer NHD)", "1", "(N, 2, B, H, D)", "同上", "dim=1: K/V 在每个 block 内交织\nstride_order 控制物理排布 (NHD vs HND)\n_reshape_kv_cache() 通过 permute 切换"],
    ["MLA bf16\n(DS V3.x)", "1", "(N, B, head_size)\nhead_size=576", "head_size\n×2=1152B", "同一行内: kv_lora[:512] + k_rope[512:]\nkernel 通过 offset 在最后维切分\n→ host 侧只看到 flat tensor, 不感知内部拆分"],
    ["MLA fp8\n(DS V3.2)", "1", "(N, B, 656)", "656 B", "字节级自定义布局:\n  kv_lora: 512B fp8 | k_rope: 128B fp8\n  metadata/align: 16B\nkernel 按 struct 解析字节 → host 侧透明"],
    ["MLA fp8\n(DS V4)", "1", "(N, B/C, 584)", "584B\n/C", "压缩布局: NoPE 448B + RoPE 128B\n  + fp8 scale 8B = 584B/token\nstorage_block_size = B // compress_ratio"],
    ["Mamba\n(GDN/Jamba)", "1 buffer\n+N views", "N 个 as_strided\nview", "sum(shapes\n×dtype)", "一个 int8 buffer 包含多个 state tensor\n每个 tensor 通过 as_strided 映射:\n  stride[0] = page_size // dtype_size (跨页)\n  storage_offset = 前序 tensor 的末尾字节"],
    ["NVFP4\nGQA", "1", "(N, 2, B, H, F)\nF=D/2+D/16", "H×F\n(F<D)", "fp4 数据 + fp8 block scale 交织在 full_dim\n每 16 个 fp4 元素配 1 个 fp8 scale\n→ 存储量减半但精度可接受"],
    ["FP8 PTH\n(Per-Token-Head)", "1 buffer\n+2 scale区", "(2, N, B, H, D)\n+ 2×(N,B,H)×f32", "同GQA\n+scales", "K/V 数据区 + 独立 K Scales 区 + V Scales 区\n3 块连续区域在同一 int8 buffer 内\nScales 字节已在 page_size_bytes 中预留"],
]
make_table(s, Inches(0.15), Inches(1.05), Inches(13.0), Inches(4.1), layout_full, header_color='188C30')

# Visual: how different data types coexist in one buffer
arch_box(s, Inches(0.3), Inches(5.35), Inches(4.0), [
    "MLA bf16 单 tensor: (N, B, 576)",
    "┌──────────────────────────────┐",
    "│  kv_lora (512) │ k_rope (64) │",
    "│  潜在向量       │ 位置编码     │",
    "└──────────────────────────────┘",
    "kernel: k_nope=tensor[...,:512]",
    "        k_rope=tensor[...,512:]",
    "→ host 侧 view(bf16).view(N,B,576) = done",
])
arch_box(s, Inches(4.5), Inches(5.35), Inches(4.2), [
    "MLA fp8 (DS V3.2): (N, B, 656)",
    "┌────────────────────────────────────┐",
    "│ kv_lora 512B │ k_rope 128B │ align │",
    "│   (fp8)      │   (fp8×2)   │  16B  │",
    "└────────────────────────────────────┘",
    "kernel: read as uint8 blob → parse struct",
    "→ host 侧 view(uint8).view(N,B,656) = done",
])
arch_box(s, Inches(8.9), Inches(5.35), Inches(4.1), [
    "Mamba: 1 buffer → N as_strided views",
    "┌──────────────────────────────────┐",
    "│conv_state(N,3,8192)│ssm(N,32,128,128)│",
    "│     float32         │    float32      │",
    "└──────────────────────────────────┘",
    "每个 view: as_strided(raw.view(f32),",
    "  size=target, stride=(page_elems, ...),",
    "  storage_offset=prev_end)",
])

# ═══════════════════════════════════════════
# PART 2: Spec 层
# ═══════════════════════════════════════════
divider_slide(9, 30, "PART 2 · Spec 层", "Spec 创建机制对比",
              "谁创建 Spec · 语义范围多大 · 以 DS V3.2 为例逐代码对比")

# SLIDE 10 — Spec 层：上游 vs Ascend 对比
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_PRIMARY); slide_num(s, 10, 30)
tb(s, Inches(0.6), Inches(0.25), Inches(12), Inches(0.5), "Spec 层：上游「类型即语义」vs Ascend「万能类」", sz=Pt(24), color=BLUE, bold=True)

card(s, Inches(0.3), Inches(0.85), Inches(6.0), Inches(2.8),
     "上游: KVCacheSpec 层次体系 (类型即语义)", [
         "KVCacheSpec → AttentionSpec → FullAttentionSpec",
         "  → MLAAttentionSpec (cache_dtype_str, compress_ratio)",
         "  → SlidingWindowSpec, ChunkedLocalAttentionSpec...",
         "每种 attention 类型 → 一种 Spec 子类",
         "Frozen dataclass: 不可变，声明式",
         "page_size_bytes: 每个子类独立实现，无分支",
         "merge(): 同类型层共享 block table 的安全校验",
     ], tc=GREEN, accent=GREEN)

card(s, Inches(6.6), Inches(0.85), Inches(6.1), Inches(2.8),
     "Ascend: AscendMLAAttentionSpec「万能类」", [
         "所有 MLA 变体挤进同一个 AscendMLAAttentionSpec",
         "新增 6 个 NPU 特有字段:",
         "  scale_dim, scale_dtype — DS V4 indexer scale",
         "  sparse_head_dim — DS V3.2/GLM5.1 三元素 tuple",
         "  cache_sparse_c8 — A3 int8 / A5 fp8 开关",
         "  c8_k_cache_dtype, c8_k_scale_cache_dtype",
         "→ 类型擦除: isinstance()无区分力",
         "→ 下游 model_runner 被迫 if-else 恢复类型信息",
     ], tc=RED, accent=RED)

code_box(s, Inches(0.3), Inches(3.9), Inches(6.0), [
    "# 上游: DS V3.2 (flash_mla_sparse.py)",
    "class FlashMLASparseBackend:",
    "    @staticmethod",
    "    def get_kv_cache_spec(...):",
    "        return MLAAttentionSpec(",
    "            block_size=block_size,",
    "            num_kv_heads=1,",
    "            head_size=576,",
    "            dtype=bf16,",
    "            cache_dtype_str='auto',",
    "        )",
    "# 只有 1 个 Spec 类型, 7 个字段, 无 NPU 概念",
], sz=Pt(7))

code_box(s, Inches(6.6), Inches(3.9), Inches(6.1), [
    "# Ascend: DS V3.2 (model_runner_v1.py:4855)",
    "class NPUModelRunner:",
    "    def get_kv_cache_spec(self, layer_name):",
    "        ...",
    "        if self.use_sparse:",
    "            return AscendMLAAttentionSpec(",
    "                block_size=self.block_size,",
    "                num_kv_heads=1,",
    "                head_size=sum(self.sparse_head_dim),",
    "                sparse_head_dim=(512, 64, 128),",
    "                dtype=bf16,",
    "                cache_dtype_str=...",
    "                cache_sparse_c8=is_c8_layer(layer),",
    "                c8_k_cache_dtype=int8|float8_e4m3fn,",
    "                c8_k_scale_cache_dtype=fp16|fp32,",
    "            )",
    "# 12+ 字段, cache_sparse_c8 决定后续所有分支",
], sz=Pt(7))

callout(s, Inches(0.3), Inches(5.9), Inches(12.4), "Spec 层核心差距",
        "上游: 类型名编码了所有语义 — MLAAttentionSpec vs SlidingWindowMLASpec 一眼就看出差异。Ascend: 所有 MLA 变体都是 AscendMLAAttentionSpec，只能通过字段值 (cache_sparse_c8? use_compress?) 间接推断行为——Spec 层本身没有解释自己的能力，model_runner 被迫替代 Spec 做解释。", AMBER)

# ═══════════════════════════════════════════
# PART 3: Backend 层
# ═══════════════════════════════════════════
divider_slide(11, 30, "PART 3 · Backend 层", "AttentionBackend 分发与 get_kv_cache_shape() 接口",
              "4 类 Ascend Backend vs 上游 Backend · 上游返回完整 shape vs Ascend 返回「参考 shape」")

# SLIDE 12 — Backend 层：以 DS V3.2 Sparse MLA 为例
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_PRIMARY); slide_num(s, 12, 30)
tb(s, Inches(0.6), Inches(0.25), Inches(12), Inches(0.5), "Backend 层：get_kv_cache_shape() — 一个返回值的信息量差距", sz=Pt(24), color=BLUE, bold=True)

code_box(s, Inches(0.3), Inches(0.9), Inches(6.0), [
    "# 上游: FlashMLASparseBackend (flashmla_sparse.py)",
    "@staticmethod",
    "def get_kv_cache_shape(",
    "    num_blocks, block_size, num_kv_heads,",
    "    head_size, cache_dtype_str='auto'",
    ") -> tuple[int, ...]:",
    "    if cache_dtype_str == 'fp8_ds_mla':",
    "        return (num_blocks, block_size, 656)",
    "    else:",
    "        return (num_blocks, block_size, head_size)",
    "",
    "# 返回 (N, B, 576) — 完整 shape, ready-to-use",
    "# model_runner: tensor.view(bf16).view(N,B,576) = OK",
], sz=Pt(7))

code_box(s, Inches(6.6), Inches(0.9), Inches(6.1), [
    "# Ascend: AscendSFABackend (sfa_v1.py:120)",
    "@staticmethod",
    "def get_kv_cache_shape(",
    "    num_blocks, block_size, num_kv_heads,",
    "    head_size, cache_dtype_str=''",
    ") -> tuple[int, ...]:",
    "    return (num_blocks, block_size,",
    "            num_kv_heads, head_size)",
    "",
    "# 返回 (N, B, 1, 576) — 「参考 shape」!",
    "# 实际需要 3~4 tensor:",
    "#   k_nope(N,B,1,512), k_rope(N,B,1,64),",
    "#   dsa_k(N,B,1,128), [dsa_k_scale(N,B,1,1)]",
    "# → 后端只回答了 1/N! 剩下靠在 model_runner 里查 layer",
], sz=Pt(7))

# The 1/N information problem visual
card(s, Inches(0.3), Inches(4.1), Inches(6.0), Inches(2.7),
     "上游 Backend: 返回即就绪", [
         "FlashAttentionBackend → (2, N, B, H, D)",
         "FlashMLABackend → (N, B, 576)",
         "FlashMLASparseBackend → (N, B, 576|656)",
         "Mamba backend → 通过 MambaSpec 的多 shapes/dtypes",
         "一个 shape = 一个 tensor 的全部信息",
         "model_runner: 拿到 shape → view → done",
         "→ model_runner ~60 行, 0 分支",
     ], tc=GREEN, accent=GREEN)

card(s, Inches(6.6), Inches(4.1), Inches(6.1), Inches(2.7),
     "Ascend Backend: 每个都只返回「参考 shape」", [
         "AscendAttentionBackend → (2, N, B, H, D) ← 这个还好",
         "AscendMLABackend → (N, B, 1, 576) ← 缺 512+64 拆分",
         "AscendSFABackend → (N, B, 1, 576) ← 缺 3~4 tensor 信息",
         "AscendDSABackend → (N, B, 1, head_size) ← 缺 scale",
         "model_runner 被迫: 查 layer 属性 → 算 split ratio",
         "  → 拆 tensor → 算各自 shape → view → 组装 tuple",
         "→ model_runner ~600 行, 15+ 分支",
     ], tc=RED, accent=RED)

# ═══════════════════════════════════════════
# PART 4: Allocate & Reshape 层
# ═══════════════════════════════════════════
divider_slide(13, 30, "PART 4 · Allocate & Reshape 层", "分配与重塑管线：上游 30 行 vs Ascend 185 行",
              "DS V3.2 Sparse MLA 完整代码拆解 · DS V4 as_strided overlay")

# SLIDE 14 — Allocate 管线对比
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_PRIMARY); slide_num(s, 14, 30)
tb(s, Inches(0.6), Inches(0.25), Inches(12), Inches(0.5), "Allocate 管线：上游 30 行 vs Ascend 185 行", sz=Pt(24), color=AMBER, bold=True)

code_box(s, Inches(0.3), Inches(0.9), Inches(6.0), [
    "# 上游 vLLM (attn_utils.py:129-143)",
    "def _allocate_kv_cache(kv_cache_config, device):",
    "    kv_cache_raw_tensors = {}",
    "    for tensor in kv_cache_config.kv_cache_tensors:",
    "        t = torch.zeros(tensor.size,",
    "                       dtype=torch.int8,",
    "                       device=device)",
    "        for layer_name in tensor.shared_by:",
    "            kv_cache_raw_tensors[layer_name] = t",
    "    return kv_cache_raw_tensors",
    "",
    "# 结束。30 行。0 分支。所有模型通用。",
], sz=Pt(7))

code_box(s, Inches(6.6), Inches(0.9), Inches(6.1), [
    "# Ascend (model_runner_v1.py:4055-4236)",
    "def _allocate_kv_cache_tensors(self, ...):",
    "    for kv_cache_tensor in kv_cache_config:",
    "        if is_hybrid_or_mamba:        # 分支1",
    "            ...  # 单 tensor, stride 适配",
    "        elif use_compress:            # 分支2",
    "            ...  # DS V4: 单 int8 buffer",
    "        else:                         # 分支3",
    "            if use_sparse:            # 分支3a",
    "                if C8 and is_A5:      # 分支3a-i",
    "                    # 3 tensor: ckv+dsa_k+scale",
    "                elif C8:             # 分支3a-ii",
    "                    # 4 tensor: k+v+dsa_k+scale",
    "                else:                # 分支3a-iii",
    "                    # 3 tensor: k+v+dsa_k",
    "            else:                     # 分支3b",
    "                # 2 tensor: k+v (split by ratio)",
    "# 185 行。决策树深度 3。新增模型: 改 3~5 处。",
], sz=Pt(7))

# Decision tree visualization
card(s, Inches(0.3), Inches(4.1), Inches(12.4), Inches(2.8),
     "Ascend Allocate 决策树：每次新增模型类型都要在此添加分支", [
         "layer_type?",
         "  ├─ hybrid/mamba → 单 tensor + stride 适配",
         "  ├─ compress (DS V4) → 单 int8 buffer (reshape 时 as_strided)",
         "  └─ standard attention",
         "       ├─ use_sparse? (DS V3.2 / GLM5.1)",
         "       │    ├─ C8 + A5?  → 3 tensor: ckv(fp8) + dsa_k(fp8) + scale(fp32)",
         "       │    ├─ C8 + A3?  → 4 tensor: k(bf16) + v(bf16) + dsa_k(int8) + scale(fp16)",
         "       │    └─ no C8     → 3 tensor: k(bf16) + v(bf16) + dsa_k(bf16)",
         "       └─ standard / MLA",
         "            ├─ is_mla? → 查 layer 获取 kv_lora_rank + qk_rope_head_dim → 按比例 split",
         "            └─ GQA     → 按 kv_head_dim 比例 split",
         "",
         "核心问题: 决策逻辑本应由 Backend/Spec 封装，却散落在 model_runner 的 if-else 中。",
     ], tc=RED, accent=RED)

# ═══════════════════════════════════════════
# SLIDE 15 — 实战拆解：DS V3.2 Sparse MLA 完整链路 (重构版)
# ═══════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_PRIMARY); slide_num(s, 15, 30)
tb(s, Inches(0.6), Inches(0.15), Inches(12), Inches(0.35),
   "实战拆解：DS V3.2 Sparse MLA — 从 Spec 到 Reshape 的完整链路",
   sz=Pt(22), color=RED, bold=True)
tb(s, Inches(0.6), Inches(0.48), Inches(12), Inches(0.22),
   "追踪一个模型经过 4 层设计 → 每一步都在丢失信息 → 最终 model_runner 被迫承载全部复杂性",
   sz=Pt(10), color=TEXT_SEC)

# ── LEFT COLUMN: 4-stage pipeline with explicit loss annotations ──
# Stage ①: Spec
code_box(s, Inches(0.3), Inches(0.82), Inches(7.2), [
    "① Spec 层: AscendMLAAttentionSpec (万能类)",
    "  sparse_head_dim=(512,64,128)  cache_sparse_c8=True/False",
    "  scale_dim / scale_dtype / c8_k_cache_dtype / c8_k_scale_cache_dtype",
    "  ⚠ 信息丢失: 类型擦除 — isinstance(spec) 不区分 5 种 MLA 变体",
    "    5 种模型共用同一个类，下游靠 if cache_sparse_c8 恢复类型信息",
], sz=Pt(7))
tb(s, Inches(3.6), Inches(1.65), Inches(0.5), Inches(0.18), "▼", sz=Pt(9), color=RED, align=PP_ALIGN.CENTER)

# Stage ②: Backend
code_box(s, Inches(0.3), Inches(1.88), Inches(7.2), [
    "② Backend 层: AscendSFABackend.get_kv_cache_shape()",
    "  return (N, B, 1, 576)  ← 只能返回 1/N 信息: 一个\"参考 shape\"",
    "  实际: k_nope(N,B,1,512) + k_rope(N,B,1,64) + dsa_k(N,B,1,128) + dsa_scale(N,B,1,?)",
    "  ⚠ 信息丢失: multi-tensor shape — 接口假设\"1 attention = 1 tensor\"",
    "    NPU 上拆成 3~4 tensor 分别存储，但接口只返回其中 1 个的等价 shape",
], sz=Pt(7))
tb(s, Inches(3.6), Inches(2.71), Inches(0.5), Inches(0.18), "▼", sz=Pt(9), color=RED, align=PP_ALIGN.CENTER)

# Stage ③: Allocate
code_box(s, Inches(0.3), Inches(2.94), Inches(7.2), [
    "③ Allocate 层: model_runner 被迫绕过 Backend",
    "  kv_lora_rank = attn_layer.kv_lora_rank         ← 直连 Attention layer 掏参数",
    "  qk_rope_head_dim = attn_layer.qk_rope_head_dim  ← Backend 的 shape 不够用",
    "  ⚠ 信息丢失: 分层抽象 — Backend 形同虚设，model_runner 越过接口直取内部状态",
    "    C8+A5 → 3 tensor (ckv+dsa_k+scale) | C8+A3 → 4 tensor | no C8 → 3 tensor",
], sz=Pt(7))
tb(s, Inches(3.6), Inches(3.77), Inches(0.5), Inches(0.18), "▼", sz=Pt(9), color=RED, align=PP_ALIGN.CENTER)

# Stage ④: Reshape
code_box(s, Inches(0.3), Inches(4.00), Inches(7.2), [
    "④ Reshape 层: _reshape_kv_cache_tensors() — 分支爆炸",
    "  A5 C8: k_shape=(N,B,1, kv_lora+rope*2+4*4)  ckv merged, v_cache=None",
    "  A3 C8: k_shape=(N,B,1,512) v_shape=(N,B,1,64)  k/v 分离",
    "  no C8: k_shape=(N,B,1,512) v_shape=(N,B,1,64)  bf16 全精度",
    "  ⚠ 信息丢失: 代码可维护性 — C8 × A5/A3 × device 三维分支，~350 行，15+ if-else",
], sz=Pt(7))

# ── RIGHT COLUMN: Upstream reference + cumulative consequence ──
code_box(s, Inches(7.8), Inches(0.82), Inches(5.2), [
    "上游 DS V3.2 Sparse: 同一条链路，正确形态",
    "",
    "FlashMLASparseBackend.get_kv_cache_spec():",
    "  return MLAAttentionSpec(         ← 1 类型 7 字段",
    "    block_size=..., num_kv_heads=1,",
    "    head_size=576, dtype=bf16,     ← 零 NPU 概念",
    "    cache_dtype_str=\"auto\")",
    "",
    "FlashMLASparseBackend.get_kv_cache_shape():",
    "  return (N, B, 576)               ← 完整 shape",
    "",
    "model_runner reshape:",
    "  tensor.view(bf16).view(N,B,576)  ← ~30 行，无分支",
], sz=Pt(7))

# Cumulative consequence
card(s, Inches(7.8), Inches(3.20), Inches(5.2), Inches(1.95),
     "累积后果: model_runner 承载了三层的全部职责", [
         "Spec 层的类型区分  →  model_runner 里 if-else 判断 C8/非C8",
         "Layout 层的存储计算 →  model_runner 里 page_size_bytes 三分支",
         "Backend 层的 shape →  model_runner 绕过 Backend 直连 attn_layer",
         "",
         "get_kv_cache_spec()        ~200 行",
         "_reshape_kv_cache_tensors() ~350 行",
         "= ~550 行 model_runner 代码, 15+ if-else 分支",
     ], tc=RED, accent=RED)

# Causal chain callout
callout(s, Inches(7.8), Inches(5.30), Inches(5.2),
        "根因链",
        "Spec 万能类 (类型擦除) → Backend 接口不足 (信息量 1/N)\n→ Layout 缺失 (存储计算无承载层) → 三者叠加\n→ model_runner 作为最终执行者被迫接管所有上游放弃的职责\n→ 这不是\"方法写太长\"的风格问题，是三层设计缺陷的必然结果",
        accent=RED)

# ── BOTTOM: 3 tensor layout comparison (compact) ──
card(s, Inches(0.3), Inches(5.15), Inches(2.3), Inches(2.1),
     "bf16 (no C8): 3 tensor", [
         "k_nope  (N,B,1,512) bf16",
         "rope    (N,B,1,64)  bf16",
         "dsa_k   (N,B,1,128) bf16",
         "──────────────",
         "Per token: 1408 B",
     ], tc=BLUE, accent=BLUE)

card(s, Inches(2.8), Inches(5.15), Inches(2.3), Inches(2.1),
     "int8 C8 (A3): 4 tensor", [
         "k_nope  (N,B,1,512) bf16",
         "rope    (N,B,1,64)  bf16",
         "dsa_k   (N,B,1,128) int8",
         "scale   (N,B,1,1)   fp16",
         "──────────────",
         "Per token: 1282 B",
     ], tc=AMBER, accent=AMBER)

card(s, Inches(5.3), Inches(5.15), Inches(2.3), Inches(2.1),
     "fp8 C8 (A5): 3 tensor", [
         "ckv     (N,B,1,704) fp8",
         "  ← kv_lora + rope 合并!",
         "dsa_k   (N,B,1,128) fp8",
         "scale   (N,B,1,1)   fp32",
         "──────────────",
         "Per token: 644 B",
     ], tc=RED, accent=RED)

# ═══════════════════════════════════════════
# SLIDE 16 — Reshape 管线对比：上游 vs Ascend (重构版)
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_PRIMARY); slide_num(s, 16, 30)
tb(s, Inches(0.6), Inches(0.15), Inches(12), Inches(0.35),
   "Reshape 管线：上游 ~30 行 vs Ascend ~350 行 — as_strided 到底在服务什么？",
   sz=Pt(22), color=AMBER, bold=True)
tb(s, Inches(0.6), Inches(0.48), Inches(12), Inches(0.22),
   "同一个工具 (as_strided): 上游服务通用运行时 (page padding) → Ascend 被迫承载模型语义 (compress/overlay/epilog) → 职责错配",
   sz=Pt(10), color=TEXT_SEC)

# ── LEFT COLUMN: Three-tier descent into complexity ──
# Tier ①: Upstream reshape
code_box(s, Inches(0.3), Inches(0.82), Inches(7.2), [
    "① 上游 Reshape: as_strided 仅服务 page padding (通用需求，与模型无关)",
    "  def _reshape_kv_cache(config, raw, backends):",
    "      for group in config.kv_cache_groups:",
    "          kv_tensor = raw[layer].view(dtype)",
    "          if spec.page_size_padded:      ← 唯一分支条件: 运行时 page padding",
    "              kv_cache = as_strided(kv_tensor, shape, stride)",
    "          else:",
    "              kv_cache = kv_tensor.view(shape)",
    "  # ~30 行。所有模型类型走同一条路径。as_strided 的 stride 由 config 决定，不是模型。",
], sz=Pt(7))
tb(s, Inches(3.6), Inches(1.70), Inches(0.5), Inches(0.18), "▼", sz=Pt(9), color=AMBER, align=PP_ALIGN.CENTER)

# Tier ②: Ascend reshape — model-type branching
code_box(s, Inches(0.3), Inches(1.93), Inches(7.2), [
    "② Ascend Reshape: _reshape_kv_cache_tensors() — 每个模型类型一个分支",
    "  if use_compress:              # DS V4 Compress MLA",
    "      A5: 3 as_strided views (K + scale + K+scale overlay) ← epilog kernel 需要叠加",
    "      A3: 2 as_strided views (K + scale)                    ← 不同硬件不同视图数",
    "  elif use_sparse:             # DS V3.2 / GLM5.1 Sparse MLA",
    "      C8+A5: CKV merged (v_cache=None) | C8+A3: k/v 分离 | no C8: bf16 全精度",
    "  elif hybrid:                 # Qwen3.5: attention + mamba 共存，切片适配",
    "  elif cache_only: ...         # 纯 cache 场景",
    "  else: ...                    # GQA / 标准 MLA",
    "  # ~350 行。as_strided 的 stride 由 模型架构(NPU算子需求) 决定，不是通用参数。",
], sz=Pt(7))
tb(s, Inches(3.6), Inches(3.38), Inches(0.5), Inches(0.18), "▼", sz=Pt(9), color=AMBER, align=PP_ALIGN.CENTER)

# Tier ③: DS V4 — the paradox
card(s, Inches(0.3), Inches(3.61), Inches(7.2), Inches(1.65),
     "③ 极端案例: DS V4 Compress MLA — 最简单的 Allocate，最复杂的 Reshape", [
         "Allocate: 1 个 int8 buffer，不拆分 ← 物理上最简单",
         "Reshape: 3 个 as_strided overlay views ← 逻辑上最复杂",
         "  View 1: K (独立 tensor)     View 2: scale (独立 tensor)",
         "  View 3: K + scale 叠加视图 (offset=0 in raw buffer, 不同 stride)",
         "  ← epilog kernel 需要同一个 buffer 上的叠加访问，这是 NPU 算子的语义需求",
         "compress_ratio C4/C128 → storage_block_size = block_size // compress_ratio → 时间维度压缩",
     ], tc=PURPLE, accent=PURPLE)

# Punchline
callout(s, Inches(0.3), Inches(5.45), Inches(7.2),
        "复杂度倒挂的根因",
        "Allocate 简单 = 模型本身简单 (1 种 tensor，不拆分)\nReshape 复杂 = NPU 算子要求复杂的视图叠加 (epilog kernel 的语义下沉到了 reshape)\n→ 上游 Spec/Layout 没有告诉 reshape \"这块 buffer 怎么解读\"\n→ reshape 只能靠自己推导 compress ratio / overlay stride / scale dtype\n→ DS V4 不是\"特例\"，是 reshape 承载模型语义的必然结果",
        accent=PURPLE)

# ── RIGHT COLUMN: as_strided purpose comparison ──
card(s, Inches(7.8), Inches(0.82), Inches(5.2), Inches(2.65),
     "as_strided 用途对比: 同一个工具，截然不同的语义层级", [
         "上游 as_strided:",
         "  触发条件: spec.page_size_padded (运行时通用参数)",
         "  决定者:   KV cache config (page size 配置)",
         "  语义层级: 通用运行时 (与模型架构、硬件平台无关)",
         "  适用范围: 所有模型类型共用同一逻辑",
         "",
         "Ascend as_strided:",
         "  触发条件: use_compress / use_sparse / hybrid / ... (模型类型)",
         "  决定者:   模型架构 (compress_ratio) + NPU 算子需求 (epilog kernel)",
         "  语义层级: 模型 + 硬件 特定 (换一个模型或 NPU 架构就要改)",
         "  适用范围: 每个模型类型一个独立分支，无法复用",
     ], tc=AMBER, accent=AMBER)

# Cumulative consequence
card(s, Inches(7.8), Inches(3.68), Inches(5.2), Inches(1.55),
     "Reshape 复杂度从哪里来？", [
         "上游 reshape 的 as_strided 条件: 1 个 (page_size_padded)",
         "Ascend reshape 的 as_strided 条件: 5+ (compress/sparse/hybrid/...)",
         "  × 每个条件内部还有二次分支: A5/A3, C8/非C8, GQA/MLA hybrid",
         "  = 15+ 个独立 reshape 路径",
         "",
         "这不是 \"代码没写好\"。是上游三层 (Spec/Backend/Layout) 没有提供",
         "\"这个 buffer 怎么解读\" 的信息，reshape 只能自己推导。",
     ], tc=RED, accent=RED)

# Bottom right: hybrid note (compact)
card(s, Inches(7.8), Inches(5.45), Inches(5.2), Inches(1.65),
     "另一个维度: Hybrid (Qwen3.5) — 两种 attention 共存", [
         "FullAttention + Mamba 混合 → 单 tensor 分配 → reshape 切片适配",
         "GQA hybrid: padding strip → K region [:attn_page_size] + V region [attn_page_size:]",
         "MLA hybrid: nope_page + rope_page + padding → slice back to views",
         "→ 这不是 reshape 的职责，但上游没有提供 region 映射信息",
     ], tc=CYAN, accent=CYAN)

# ═══════════════════════════════════════════
# PART 5: 按模型类型的完整差异链路
# ═══════════════════════════════════════════
divider_slide(17, 30, "PART 5 · 模型类型差异", "五种模型类型的完整链路对比",
              "从 GQA → 标准 MLA → Sparse MLA → Compress MLA → Hybrid 差异递增")

# SLIDE 18 — 跨模型对比总表 (重构版: 上游 vs Ascend 逐列对比)
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_PRIMARY); slide_num(s, 18, 30)
tb(s, Inches(0.6), Inches(0.15), Inches(12), Inches(0.35),
   "五种模型类型：上游 vs Ascend — Spec → Backend → Allocate → Reshape 逐列对比",
   sz=Pt(22), color=BLUE, bold=True)
tb(s, Inches(0.6), Inches(0.48), Inches(12), Inches(0.22),
   "每一列同时展示上游基线 (绿色) 和 Ascend 现状 (橙色) → 差距 (Δ) 随模型复杂度递增",
   sz=Pt(10), color=TEXT_SEC)

model_data = [
    ["模型类型", "Spec\n上游 → Ascend", "Backend\n上游 → Ascend", "Allocate (tensor数)\n上游 → Ascend", "Reshape 方式\n上游 → Ascend", "差距 (Δ)"],
    ["GQA\n(Qwen3 MoE)",
     "FullAttentionSpec\n→ 同上游 ✓",
     "FlashAttention\n→ AscendAttention",
     "1 连续 tensor\n→ 2 tensor (K/V物理拆分)",
     "view(N,B,H,D)\n→ .view() 同上游",
     "K/V 物理拆分\n(PD分离需求)\n差异: ★"],
    ["标准 MLA\n(DS V3.1)",
     "MLAAttentionSpec\n→ AscendMLAAttentionSpec\n(+scale_dim/dtype)",
     "FlashMLA\n→ AscendMLA",
     "1 连续 tensor\n→ 2 tensor (K/V分离)",
     "view(N,B,576)\n→ view + 查 layer 属性\n(head_size 丢失512+64)",
     "head_size=576\n丢失 512+64 拆分\n差异: ★★"],
    ["Sparse MLA\n(DS V3.2/GLM5.1)",
     "MLAAttentionSpec\n→ AscendMLAAttentionSpec\n(+sparse_head_dim/C8)",
     "FlashMLASparse\n→ AscendSFA",
     "1 连续 tensor\n→ 3~4 tensor (C8/A5分支)",
     "view(N,B,576)\n→ unpack + view\n+ C8×A5/A3 多分支",
     "indexer + C8 量化\nA5/A3 硬件分支\n差异: ★★★"],
    ["Compress MLA\n(DS V4)",
     "MLAAttentionSpec\n→ AscendMLAAttentionSpec\n(+compress_ratio/scale)",
     "FlashMLASparse\n→ AscendDSA",
     "1 fp8 tensor\n→ 1 int8 tensor",
     "view(N,B,656)\n→ as_strided overlay\n(3 views, epilog kernel)",
     "时间压缩 + fp8\nas_strided 承载模型语义\n差异: ★★★★"],
    ["Hybrid\n(Qwen3.5)",
     "FullAttn + MambaSpec\n→ 3种 spec 共存",
     "FlashAttn + GDN\n→ 混合 Backend",
     "8 shared tensors\n→ 1 shared tensor",
     "view + as_strided\n→ slice + strip + padding\n(region 映射手动推导)",
     "attn + mamba 共存\nlayout 互斥, 手动切片\n差异: ★★★★★"],
]
make_table(s, Inches(0.3), Inches(0.82), Inches(12.4), Inches(3.3), model_data)

# ── Two summary cards: upstream baseline vs Ascend pattern ──
card(s, Inches(0.3), Inches(4.35), Inches(6.2), Inches(1.45),
     "上游基线: 5 种模型走同一条简洁路径", [
         "Spec:   每种模型 1 个专用 Spec 子类, isinstance() 编码全部语义",
         "Backend: 1 个静态方法, get_kv_cache_shape() 返回完整可用 shape",
         "Allocate: 连续单 tensor 分配, 无拆分逻辑",
         "Reshape: view(N,B,H) 一行完成, as_strided 仅用于 page padding",
         "→ 整个 pipeline ~100 行, model_runner 不感知模型差异",
     ], tc=GREEN, accent=GREEN)

card(s, Inches(6.8), Inches(4.35), Inches(6.2), Inches(1.45),
     "Ascend 现状: 每种模型一个独立分支, 复杂度递增", [
         "Spec:   5 种变体共用万能类, isinstance() 无区分力 → if-else 恢复类型",
         "Backend: get_kv_cache_shape() 返回 1/N 信息 → model_runner 绕过接口",
         "Allocate: 2~4 tensor, C8×A5/A3×device 三维分支",
         "Reshape: unpack / as_strided overlay / slice+strip 每个模型独立逻辑",
         "→ model_runner ~600 行, 15+ if-else, 每加一种模型改一次",
     ], tc=RED, accent=RED)

# Punchline
callout(s, Inches(0.3), Inches(5.98), Inches(12.4),
        "关键洞察",
        "GQA (差异 ★) → 标准 MLA (★★) → Sparse MLA (★★★) → Compress MLA (★★★★) → Hybrid (★★★★★)\n差距不是线性叠加——每增加一个模型特性 (sparse/C8/compress/hybrid)，Ascend pipeline 就多一层 if-else\n上游通过 Spec 子类 + Layout 多态吸收模型差异；Ascend 把这部分职责全部下沉到了 model_runner",
        accent=RED)

# ═══════════════════════════════════════════
# PART 5.5: Qwen3.5 精选案例
# ═══════════════════════════════════════════
divider_slide(19, 30, "PART 5.5 · Qwen3.5 Hybrid 精选案例", "以 Qwen3.5 为例：从 Spec 到 Reshape 的完整链路拆解",
              "32 层 (8 full_attn + 24 linear_attn) · 上游 vs Ascend 代码级逐行对比")

# SLIDE 20 — Qwen3.5: 上游如何用类型系统优雅处理混合架构
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_PRIMARY); slide_num(s, 20, 30)
tb(s, Inches(0.6), Inches(0.15), Inches(12), Inches(0.35),
   "Qwen3.5 Hybrid — 上游如何用类型系统处理混合架构 (FullAttention + Mamba)",
   sz=Pt(22), color=GREEN, bold=True)
tb(s, Inches(0.6), Inches(0.48), Inches(12), Inches(0.22),
   "核心原则: model_runner 从不检查 layer_type — Spec 子类 + Backend 多态吸收所有模型差异",
   sz=Pt(10), color=TEXT_SEC)

# ── Row 1: Architecture + Spec (two cards) ──
card(s, Inches(0.3), Inches(0.82), Inches(4.0), Inches(2.15),
     "Qwen3.5 架构: 两种 attention 共存", [
         "32 层 (dense): 每第 4 层 full_attention (8层)",
         "  其余 24 层 linear_attention (GatedDeltaNet)",
         "full_attn: num_kv_heads=4, head_dim=256",
         "linear_attn: num_k_heads=16, num_v_heads=32",
         "  conv_state=(3,8192), ssm_state=(32,128,128)",
         "mamba_cache_mode 强制 'align' (与 attn block 对齐)",
         "→ FullAttention + Mamba 没有公共基类!",
     ], tc=BLUE, accent=BLUE)

code_box(s, Inches(4.6), Inches(0.82), Inches(4.1), [
    "上游 Spec: 2 种原生类型，各司其职",
    "",
    "full_attn → FullAttentionSpec(",
    "  block_size=16, num_kv_heads=4,",
    "  head_size=256, dtype=bf16)",
    "  → FlashAttentionBackend    ← GQA 路径",
    "",
    "linear_attn → MambaSpec(",
    "  shapes=((3,8192),(32,128,128)),",
    "  dtypes=(f32,f32),",
    "  mamba_type='GDN_ATTN')",
    "  → GDNAttentionBackend       ← Mamba 路径",
], sz=Pt(7))

card(s, Inches(8.9), Inches(0.82), Inches(4.1), Inches(2.15),
     "上游 分组 + 共享: 8 tensor 服务 32 层", [
         "分组: 4 groups = 1 full_attn ×8 + 3 linear_attn ×8",
         "  min_num_layers=8 → group_size=8",
         "每个 KVCacheTensor 被 4 层共享:",
         "  1 attn + 3 linear → 共享同一块物理 buffer",
         "attn 使用 buffer 的 K/V 区域 (2,N,B,H,D)",
         "mamba 使用 buffer 的 conv+ssm 区域 (flatten)",
         "→ 上游不管\"内部怎么切\"— 各 Backend 自行解读",
     ], tc=CYAN, accent=CYAN)

# ── Row 2: Pipeline flow — Allocate → Reshape (two code boxes) ──
code_box(s, Inches(0.3), Inches(3.15), Inches(6.4), [
    "上游 Allocate: 30 行, 0 分支 — 不区分 attn/mamba",
    "",
    "def _allocate_kv_cache(config, device):",
    "    raw_tensors = {}",
    "    for t in config.kv_cache_tensors:     ← 统一循环",
    "        tensor = torch.zeros(t.size,",
    "                           dtype=int8,",
    "                           device=device)",
    "        for name in t.shared_by:           ← 4层共享1 buffer",
    "            raw_tensors[name] = tensor",
    "    return raw_tensors",
    "",
    "# 没有 if layer.is_attn / elif layer.is_mamba!",
    "# Spec + Backend 已经告诉了 config 该怎么分配。",
], sz=Pt(7))

code_box(s, Inches(7.0), Inches(3.15), Inches(6.0), [
    "上游 Reshape: view + as_strided_ — 零拷贝重排",
    "",
    "def _reshape_kv_cache(...):",
    "    # attn: 自然交织 K/V",
    "    shape = backend.get_shape(N,B,4,256)",
    "    # → (2, N, B, 4, 256)   ← K=idx0, V=idx1",
    "    kv_cache = raw.view(bf16).view(shape)",
    "",
    "    # hybrid layout adjust: (2,N,...) → (N,2,...)",
    "    hidden_size = kv_cache.shape[2:].numel()",
    "    kv_cache.as_strided_(",
    "        size=(N,2,B,H,D),",
    "        stride=(2*hidden_size, hidden_size, ...))",
    "    # block[i] 内: K@offset0, V@offset hidden_size",
    "    # → 与 mamba block-dim=0 对齐!",
], sz=Pt(7))

# ── Row 3: Punchline ──
callout(s, Inches(0.3), Inches(5.75), Inches(12.4),
        "上游的核心优势: 类型系统吸收差异, model_runner 零感知",
        "Spec: 两种原生类型 (FullAttentionSpec + MambaSpec) → isinstance() 区分 attn vs mamba\n"
        "Backend: FlashAttentionBackend + GDNAttentionBackend → 各自的 get_kv_cache_shape() 返回正确 shape\n"
        "Allocate: 统一循环, 30 行, 无 layer_type 判断 → config 已编码所有信息\n"
        "Reshape: view() + as_strided_() 零拷贝重排 → (2,N,...) → (N,2,...) 纯 stride 调整\n"
        "→ model_runner 从!不!检!查! layer_type — 所有模型特有逻辑封装在 Spec + Backend 内部",
        accent=GREEN)

# ═══════════════════════════════════════════
# SLIDE 21 — Qwen3.5: Ascend 管线 + 逐阶段对比
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_PRIMARY); slide_num(s, 21, 30)
tb(s, Inches(0.6), Inches(0.15), Inches(12), Inches(0.35),
   "Qwen3.5 Hybrid — Ascend 管线 + 上游 vs Ascend 逐阶段对比",
   sz=Pt(22), color=RED, bold=True)
tb(s, Inches(0.6), Inches(0.48), Inches(12), Inches(0.22),
   "Spec 复用上游 ✓ 但 Allocate/Reshape 被迫引入 hybrid 专用路径 → 差异集中在管线后半段",
   sz=Pt(10), color=TEXT_SEC)

# ── Row 1: Ascend overview (two compact cards) ──
card(s, Inches(0.3), Inches(0.82), Inches(6.4), Inches(1.55),
     "Ascend Spec + Backend: 最好的部分 (几乎无差异)", [
         "Spec: 直接复用上游 FullAttentionSpec + MambaSpec",
         "  ← Qwen3.5 不用 MLA → 唯一不需要 AscendMLAAttentionSpec 的模型!",
         "唯一修改: object.__setattr__(attn_spec, 'page_size_padded', mamba_size)",
         "  → 强制对齐到 mamba，触发 hybrid 单 tensor 路径",
         "Backend: AscendAttentionBackend (GQA) + GDNAttentionBackend",
         "  get_kv_cache_shape() → (2,N,B,4,256) — 与上游完全相同!",
         "  ← 唯一不需要「二次解读」shape 的 Ascend Backend",
     ], tc=GREEN, accent=GREEN)

card(s, Inches(7.0), Inches(0.82), Inches(6.0), Inches(1.55),
     "Ascend Allocate + Reshape: 差异集中爆发 (hybrid 专用路径)", [
         "Allocate: hybrid_with_attn_and_mamba = True",
         "  → 单 int8 tensor 分配 (表面与上游相同)",
         "  但内部 layout 完全不同 — 为后续 strip/split 预留",
         "Reshape (model_runner_v1:4370-4615, ~100 行 hybrid 专用):",
         "  Step 1: attn_page_size = prod(shape[1:]) * dtype.itemsize",
         "  Step 2: conv_padding = raw.numel() - attn_page_size * 2",
         "  Step 3: raw = raw[conv_padding:]    # strip mamba conv 前缀",
         "  Step 4: k = raw[:attn_page_size]    # K 区域 (手动计算 offset)",
         "           v = raw[attn_page_size:]   # V 区域 (手动计算 offset)",
     ], tc=RED, accent=RED)

# ── Row 2: Comparison table (centerpiece, larger) ──
qwen_data = [
    ["阶段", "上游 vLLM (基线)", "vLLM-Ascend (现状)", "差距 (Δ) 及原因"],
    ["Spec\n类型",
     "FullAttentionSpec + MambaSpec\n2 种原生类型，各司其职",
     "完全复用上游 ✓\n+ page_size_padded 强制对齐",
     "几乎无差异\n(Qwen3.5 不用 MLA)"],
    ["Backend\nshape",
     "(2,N,B,4,256)\nK/V 交织在同一 tensor",
     "(2,N,B,4,256)\n与上游完全相同!",
     "无差异\n(shape 语义足够)"],
    ["Allocate",
     "单 int8 tensor 统一循环\n30 行, 0 分支",
     "单 int8 tensor (hybrid 路径)\n内部 layout 为 strip 预留",
     "分配相同, 语义不同\n(上游不管内部布局)"],
    ["Reshape\n(核心差距)",
     "view(bf16).view(2,N,B,H,D)\n→ as_strided_ 调 stride\n零拷贝, ~5 行",
     "strip conv_pad → 手动算 K 区 offset\n→ 手动算 V 区 offset → split\n~100 行 hybrid 专用代码",
     "上游 as_strided_ 解决的事\nAscend 需 strip+split\n(NPU 不支持 stride view)"],
    ["最终\n结果",
     "attn: (N,2,B,H,D) 单 tensor\nmamba: list[tensor]\n同一 buffer 两种视图",
     "attn: (k_tensor, v_tensor) tuple\nmamba: list[tensor]\nK/V 必须物理独立",
     "K/V 必须独立 tensor\n→ 传给 NPU 算子的要求"],
]
make_table(s, Inches(0.3), Inches(2.55), Inches(12.4), Inches(2.45), qwen_data)

# ── Row 3: Buffer layout comparison (visual evidence) ──
arch_box(s, Inches(0.3), Inches(5.15), Inches(6.4), [
    "上 游 Buffer: 单 tensor, K/V 自然交织, as_strided_ 零拷贝重排",
    "",
    "  ┌───────────────────────┬───────────────────────┐",
    "  │  K block0..N (连续)    │  V block0..N (连续)    │",
    "  └───────────────────────┴───────────────────────┘",
    "  view(bf16).view(2,N,B,H,D)  →  as_strided_  →  (N,2,B,H,D)",
    "  K/V 仍在同一 tensor 内, stride 改变逻辑视图, 不物理拷贝",
])

arch_box(s, Inches(7.0), Inches(5.15), Inches(6.0), [
    "Ascend Buffer: 物理分四区域, strip + 手动 split → K/V tuple",
    "",
    "  ┌──────────┬────────────┬────────────┬──────────┐",
    "  │conv_state│ K blocks   │ V blocks   │ssm_state │",
    "  │(mamba)   │ (attn)     │ (attn)     │(mamba)   │",
    "  └──────────┴────────────┴────────────┴──────────┘",
    "   ↑ strip    ↑ raw[p:][:K]  ↑ raw[p:][K:]  ↑ sliced",
    "   conv_pad   手动计算 offset  手动计算 offset  [start:end]",
])

# ═══════════════════════════════════════════
# PART 6: 精度维度
# ═══════════════════════════════════════════
divider_slide(22, 30, "PART 6 · 精度维度", "bf16 vs fp8：与模型、硬件正交的第三维度",
              "同一模型 DS V3.2 在 bf16 / int8 C8 / fp8 C8 三种精度下的完整链路对比")

# SLIDE 23 — 精度维度总览 (重构版)
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_PRIMARY); slide_num(s, 23, 30)
tb(s, Inches(0.6), Inches(0.15), Inches(12), Inches(0.35),
   "精度维度：与模型类型、硬件设备正交的第三维度",
   sz=Pt(22), color=BLUE, bold=True)
tb(s, Inches(0.6), Inches(0.48), Inches(12), Inches(0.22),
   "精度不是模型的子属性，不是硬件的子属性——它是独立的第三维：同一模型 × 同一设备 × 不同精度 = 完全不同的 KV Cache 行为",
   sz=Pt(10), color=TEXT_SEC)

# ── Row 1: The orthogonality concept (3 cards showing 3 independent axes) ──
card(s, Inches(0.3), Inches(0.82), Inches(4.0), Inches(1.55),
     "维度 1: 模型类型 (5 种)", [
         "GQA → 标准 MLA → Sparse MLA → Compress MLA → Hybrid",
         "决定: tensor 数量, split 策略, 是否 CKV 合并",
         "上游已通过 Spec 子类吸收此维度 ✓",
     ], tc=BLUE, accent=BLUE)

card(s, Inches(4.6), Inches(0.82), Inches(4.0), Inches(1.55),
     "维度 2: 精度 (3 种)", [
         "bf16 (非 C8) → int8 C8 (A3) → fp8 C8 (A5)",
         "决定: 每个 tensor 的 dtype, scale 有无, page 字节数",
         "此维度完全独立 — 同一模型可走三种精度路径",
     ], tc=AMBER, accent=AMBER)

card(s, Inches(8.9), Inches(0.82), Inches(4.1), Inches(1.55),
     "维度 3: 硬件设备 (A3 vs A5)", [
         "A3: int8 C8, K/V 分离, 4 tensor",
         "A5: fp8 C8, CKV 合并 (v_cache=None), 3 tensor",
         "决定: CKV merged shape 公式, as_strided overlay",
     ], tc=RED, accent=RED)

# ── Row 2: DS V3.2 as the exemplar (walks through all 3 precisions) ──
tb(s, Inches(0.6), Inches(2.55), Inches(12), Inches(0.25),
   "DS V3.2 Sparse MLA — 唯一覆盖全部三种精度的模型：精度切换如何影响 pipeline 的每一层",
   sz=Pt(11), color=TEXT_PRIMARY, bold=True)

ds32_data = [
    ["pipeline 阶段", "bf16 (非 C8)", "int8 C8 (A3)", "fp8 C8 (A5)", "精度切换的影响"],
    ["Spec:\ncache_sparse_c8", "False", "True", "True", "一个 bool 翻转"],
    ["Spec:\nk_cache_dtype", "bf16", "int8", "float8_e4m3fn", "dtype 改变"],
    ["page_size_bytes\n(每 token)", "1408 B\n(512+64+128)×2", "1282 B\n512×2+64×2+128+2", "644 B\n512+128+4", "字节数 ≠ 简单线性!"],
    ["Allocate\ntensor 数", "3 tensor\nk+v+dsa_k", "4 tensor\nk+v+dsa_k+scale", "3 tensor\nckv+dsa_k+scale", "3→4→3 非单调!"],
    ["Reshape\nv_cache", "(N,B,1,64)\n独立 V tensor", "(N,B,1,64)\n独立 V tensor", "None\n合并入 CKV", "V 消失了!"],
    ["Reshape\nk_shape", "(N,B,1,512)\n独立的 K", "(N,B,1,512)\n独立的 K", "(N,B,1,704)\nCKV=K+V+pad", "K 的 shape 变了"],
]
make_table(s, Inches(0.3), Inches(2.88), Inches(12.4), Inches(2.8), ds32_data)

# ── Row 3: Punchline ──
callout(s, Inches(0.3), Inches(5.85), Inches(12.4),
        "精度的正交性意味着什么",
        "上游: Backend.get_kv_cache_shape() 内部消化精度差异 → model_runner 只看到 shape 最后一维变了 (576→656), 其余完全透明\n"
        "Ascend: 精度切换需要同时修改 5 个位置 (page_size_bytes / sparse_kv_cache_ratio / Allocate 分支 / Reshape 公式 / CKV shape)\n"
        "→ 精度不应该是 model_runner 的职责。理想方案: SparseMLALayout(precision='fp8_c8', device='A5') 在构造时消化全部精度差异",
        accent=AMBER)

# ═══════════════════════════════════════════
# SLIDE 24 — 上游 vs Ascend 精度处理对比 (重构版)
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_PRIMARY); slide_num(s, 24, 30)
tb(s, Inches(0.6), Inches(0.15), Inches(12), Inches(0.35),
   "精度处理：上游一个字段 vs Ascend 五个位置 — 差距为什么这么大？",
   sz=Pt(22), color=BLUE, bold=True)
tb(s, Inches(0.6), Inches(0.48), Inches(12), Inches(0.22),
   "根因: GPU kernel 可以内部解析打包 buffer (blob) → NPU kernel 需要 host 侧预先拆分 → 精度差异被迫泄漏到 model_runner",
   sz=Pt(10), color=TEXT_SEC)

# ── Row 1: Side-by-side code contrast ──
code_box(s, Inches(0.3), Inches(0.82), Inches(6.4), [
    "上游: 精度切换 — 只改 Backend 内部的 shape 最后一维",
    "",
    "def get_kv_cache_shape(..., cache_dtype_str='auto'):",
    "    if cache_dtype_str == 'fp8_ds_mla':",
    "        return (N, B, 656)    ← fp8 字节布局 (含 scale)",
    "    else:",
    "        return (N, B, 576)    ← bf16",
    "",
    "# model_runner 看到的: shape 从 (N,B,576) 变成 (N,B,656)",
    "# 其余一切不变 — tensor 数量不变, dtype 不变, 分支不变",
    "# → cache_dtype_str 1 个字段封装全部精度差异",
], sz=Pt(7))

code_box(s, Inches(7.0), Inches(0.82), Inches(6.0), [
    "Ascend: 精度切换 — 需要同步修改 5 个不同位置",
    "",
    "① page_size_bytes:      if C8: A5公式? A3公式? else: bf16公式",
    "② sparse_kv_cache_ratio: A5返回3-tuple, A3返回4-tuple",
    "③ Allocate 分支:         if C8+A5:3t elif C8:4t else:3t",
    "④ Reshape k_shape:       A5 CKV=(N,B,1,704), A3 k=(N,B,1,512)",
    "⑤ Reshape v_shape:       A5 v=None(CKV合并), A3 v=(N,B,1,64)",
    "",
    "# 这 5 处分散在 model_runner 不同函数里",
    "# 每加一种精度 → 需要找到所有分支点 → 手动保持一致",
], sz=Pt(7))

# ── Row 2: Root cause analysis — two contrasting cards ──
card(s, Inches(0.3), Inches(2.95), Inches(6.4), Inches(2.05),
     "上游为什么能做到 1 个字段搞定", [
         "GPU CUDA kernel 的设计模式:",
         "  接受单 uint8/int8 buffer (blob)",
         "  kernel 内部按 struct 字节布局解析各字段",
         "  → 传一个 blob 就能跑, host 侧只管总大小",
         "",
         "精度切换的影响范围: 仅 Backend.get_kv_cache_shape()",
         "  bf16 → shape 最后一维 = 576",
         "  fp8  → shape 最后一维 = 656 (含 scale 和 padding)",
         "  → model_runner 完全不感知精度变化",
     ], tc=GREEN, accent=GREEN)

card(s, Inches(7.0), Inches(2.95), Inches(6.0), Inches(2.05),
     "Ascend 为什么做不到 — NPU 算子接口的根本约束", [
         "NPU kernel 的设计模式:",
         "  不接受打包 blob, 每个子 component 必须是独立 tensor",
         "  且 dtype 必须正确 (bf16/int8/fp8/fp32 各不同)",
         "  → host 侧必须预先拆分并标注每个 tensor 的 dtype",
         "",
         "精度切换的影响范围: 遍及整个 pipeline",
         "  tensor 数量 (3→4→3), 每个 tensor 的 dtype",
         "  scale tensor 有无, CKV merge shape 公式",
         "  page_size_bytes 计算公式, sparse_kv_cache_ratio",
     ], tc=RED, accent=RED)

# ── Row 3: Punchline + fix direction ──
callout(s, Inches(0.3), Inches(5.20), Inches(12.4),
        "精度正交性的正确承载方式",
        "上游: 精度差异被 Backend 内部消化 → model_runner 只看到 shape 变化 → 1 个字段 (cache_dtype_str) 编码全部精度信息\n"
        "Ascend 当前: 精度差异散落在 5 个位置 → 每次精度切换需要手动同步 → 容易出错, 不可测试\n"
        "目标: Layout(precision, device) 构造参数接受精度+设备 → Layout 内部消化全部差异 → model_runner 回归\"不感知精度\"",
        accent=RED)

# ═══════════════════════════════════════════
# PART 7: 根因 + 对齐
# ═══════════════════════════════════════════
divider_slide(25, 30, "PART 7 · 根因 + 对齐", "差异根因分析与社区对齐路线",
              "六大根因 → Layout 重构 → 三步对齐计划")

# SLIDE 26 — 六大根因
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_PRIMARY); slide_num(s, 26, 30)
tb(s, Inches(0.6), Inches(0.25), Inches(12), Inches(0.5), "六大根因，归结为三个源头", sz=Pt(28), color=RED, bold=True)

causes = [
    ("1. K/V 物理分离", RED, "PD分离架构要求K/V独立RDMA传输。2MB对齐是硬件要求。影响所有模型类型。"),
    ("2. NPU 不支持 stride", RED, "DMA Engine只接受{base_addr,length}。GPU的SIMT LSU允许任意地址计算。根因中的根因。"),
    ("3. Spec 万能类: 类型擦除", MAGENTA, "5种MLA变体→1个AscendMLAAttentionSpec。isinstance()不编码语义，下游用if-else恢复已丢失的类型信息。"),
    ("4. Compress 交织 Layout", AMBER, "DS V4的K和scale按block交织。reshape需as_strided overlay零拷贝。A5 epilog需叠加视图。"),
    ("5. A3 vs A5 硬件差异", BLUE, "A5: CKV fp8合并(3 tensor)。A3: bf16分离(4 tensor)。device分支散落各处。"),
    ("6. Backend 接口不足", BLUE, "get_kv_cache_shape()单返回值无法表达多tensor布局。语义泄漏到model_runner。"),
]
for i, (title, clr, desc) in enumerate(causes):
    col, row = i % 3, i // 3
    x = Inches(0.3) + Inches(col*4.3); y = Inches(0.9) + Inches(row*1.7)
    card(s, x, y, Inches(4.1), Inches(1.45), title, [desc], tc=clr, accent=clr)

callout(s, Inches(0.3), Inches(4.5), Inches(12.4), "三个根本源头",
        "硬件约束: NPU算子不支持stride访问 → host侧必须预拆分tensor → 分配和重塑逻辑复杂化 (根因1,2,4,5)。"
        "Spec 设计: 类型不编码语义 → isinstance()无区分力 → 下游被迫用if-else恢复丢失的类型信息 (根因3)。"
        "接口缺陷: get_kv_cache_shape() 单返回值无法表达多tensor布局 → 信息泄漏到 model_runner (根因6)。", RED)

card(s, Inches(0.3), Inches(5.3), Inches(6.0), Inches(1.4),
     "Layout 重构：Spec 拆分 + Layout 多态", [
         "Spec 层: 拆分为4~5个子类 (AscendStandard/Sparse/C8/DSV4)",
         "  → 每个子类零 if-else，类型名=语义",
         "Layout 层: KVCacheLayout 6个子类",
         "  → spec.get_layout() 多态返回对应 Layout",
         "  → layout.allocate() / layout.reshape() 零分支",
         "结果: model_runner ~135行 + Layout 类 (独立可测)",
     ], tc=GREEN, accent=GREEN)

card(s, Inches(6.6), Inches(5.3), Inches(6.1), Inches(1.4),
     "社区对齐：5 大任务 + 三步路线", [
         "T1: KVCacheLayout 6子类 + Spec 4子类拆分 (Q3, P0)",
         "T2: 提议 get_kv_cache_layout() 为硬件无关抽象 (Q4, P0)",
         "T3: Spec字段通用化 (sparse_head_dim等) (Q4, P1)",
         "T4: 删除 model_runner patch + NPU特有字段 (Q1+, P1)",
         "T5: 算子侧 stride 支持推进 (持续, P2)",
     ], tc=BLUE, accent=BLUE)

# ═══════════════════════════════════════════
# SLIDE 27 — 差距分析 + 三步路线图 (重构版)
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_PRIMARY); slide_num(s, 27, 30)
tb(s, Inches(0.6), Inches(0.15), Inches(12), Inches(0.35),
   "差距分析：上游 vs Ascend — 三个层面 + 三步对齐",
   sz=Pt(22), color=GREEN, bold=True)
tb(s, Inches(0.6), Inches(0.48), Inches(12), Inches(0.22),
   "差距从 分组 (最小) → 分配 (关键) → 内部布局 (根本) 逐级递增 → 三步路线逐层解决",
   sz=Pt(10), color=TEXT_SEC)

# ── Row 1: Three gap levels (gradient: green → amber → red) ──
card(s, Inches(0.3), Inches(0.85), Inches(4.0), Inches(1.85),
     "层面 1: 分组 (差距 ★ — 最小)", [
         "上游: get_kv_cache_groups() 决策级联",
         "  按 page_size / block_size 自动分组",
         "Ascend: 几乎完全复用上游逻辑",
         "  仅 DS V4 有 patch (三层结构需按 page_size 分桶)",
         "对齐难度: ★ (低) — 上游已支持 DS V4",
         "  未来删除 patch 即可 → Layout 重构后自动解决",
     ], tc=GREEN, accent=GREEN)

card(s, Inches(4.6), Inches(0.85), Inches(4.0), Inches(1.85),
     "层面 2: 分配 (差距 ★★★ — 关键)", [
         "上游: 1 tensor/层, KVCacheTensor.shared_by",
         "  统一循环分配, 30 行, 0 分支",
         "Ascend: N tensors/层 (K/V 分离 + sparse 拆分)",
         "  185 行, 深度 3 决策树, 每加模型改 3~5 处",
         "对齐难度: ★★★ (中) — 需要 Layout 层封装",
         "  Layout.split_sizes() → 返回每 tensor 字节数列表",
     ], tc=AMBER, accent=AMBER)

card(s, Inches(8.9), Inches(0.85), Inches(4.1), Inches(1.85),
     "层面 3: 内部布局 (差距 ★★★★★ — 根本)", [
         "上游: kernel 内部解析字节/offset",
         "  GPU 支持 stride → as_strided 零拷贝 → host 不关心内部",
         "Ascend: host 侧预拆分 → reshape 必须知道内部结构",
         "  350 行, C8×A5/A3×precision 三维分支",
         "对齐难度: ★★★★★ (高) — 需要 Layout + 算子双管齐下",
         "  Layout.reshape() 封装 host 侧逻辑 → 算子侧长期推 stride",
     ], tc=RED, accent=RED)

# ── Row 2: Already aligned (compact) + What remains (compact) ──
card(s, Inches(0.3), Inches(2.90), Inches(6.4), Inches(1.2),
     "已对齐 (无需额外工作) ✓", [
         "✓ KVCacheSpec / KVCacheConfig / KVCacheTensor   ✓ BlockPool / SingleTypeKVCacheManager",
         "✓ KVCacheCoordinator (仅 DS V4 有 patch)    ✓ Prefix Caching hash & block 管理",
         "✓ shared_by 多层共享机制    ✓ get_kv_cache_groups() (仅 DS V4 有 patch)",
         "→ 这些属于「规划层」— 在 tensor 分配之前，与 NPU 硬件无关，可直接复用上游",
     ], tc=GREEN, accent=GREEN)

card(s, Inches(7.0), Inches(2.90), Inches(6.0), Inches(1.2),
     "待解决 (全部 P0)", [
         "✗ Spec 万能类 → 拆分为 4~5 个子类, 恢复类型即语义",
         "✗ Backend 接口不足 → 扩展为 list[tuple] 或 Layout 对象",
         "✗ Allocate 决策树 → Layout.split_sizes() 替代 185 行 if-else",
         "✗ Reshape 三维分支 → Layout.reshape() 替代 350 行分支",
         "→ 这些属于「执行层」— 直接依赖 NPU 硬件特性, 需要 Layout 层封装",
     ], tc=RED, accent=RED)

# ── Row 3: Three-phase roadmap ──
tb(s, Inches(0.6), Inches(4.35), Inches(12), Inches(0.22),
   "三步对齐路线：内部重构 → 上游 RFC → 社区合入",
   sz=Pt(11), color=TEXT_PRIMARY, bold=True)

phases = [
    ("Phase 1 (Q3):\n内部重构", GREEN,
     "目标: KVCacheLayout 6 子类 + Spec 4~5 子类拆分\n"
     "效果: model_runner ~135 行, 15+ if-else → 多态分发, 每 Layout 独立可测\n"
     "关键交付: StandardKVLayout / SparseMLALayout / CompressedMLALayout / HybridLayout / CacheOnlyLayout"),
    ("Phase 2 (Q4):\n上游 RFC", BLUE,
     "目标: 向社区提议 KVCacheLayout 为硬件无关抽象\n"
     "内容: get_kv_cache_layout() 接口 + Backend 接口扩展 + Spec 字段通用化 (sparse_head_dim 等)\n"
     "关键交付: RFC 文档 + 原型代码 + 社区讨论 → 获得上游认可"),
    ("Phase 3 (Q1+):\n社区合入 + 清理", PURPLE,
     "目标: Layout 子类合入社区 → 删除所有临时方案\n"
     "效果: 删除 model_runner patch (~600 行) → 删除 AscendMLAAttentionSpec NPU 特有字段 → 删除 patch_kv_cache_utils\n"
     "关键交付: 社区 PR 合入 + Ascend 代码清理 → 与上游基线一致"),
]
for i, (title, clr, desc) in enumerate(phases):
    y = Inches(4.65) + Inches(i*0.82)
    badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), y, Inches(2.0), Pt(40))
    r2 = max(0, min(255, int(clr[0]*0.1))); g2 = max(0, min(255, int(clr[1]*0.1))); b2 = max(0, min(255, int(clr[2]*0.1)))
    badge.fill.solid(); badge.fill.fore_color.rgb = RGBColor(r2,g2,b2); badge.line.fill.background()
    badge.adjustments[0] = 0.3
    p = badge.text_frame.paragraphs[0]; p.text = title; p.font.size = Pt(8); p.font.color.rgb = clr; p.font.bold = True
    p.font.name = 'Microsoft YaHei'; p.alignment = PP_ALIGN.CENTER
    tb(s, Inches(2.5), y, Inches(10.5), Pt(40), desc, sz=Pt(7.5), color=TEXT_SEC)

# ═══════════════════════════════════════════
# SLIDE 28 — Layout 重构方案
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_PRIMARY); slide_num(s, 28, 30)
tb(s, Inches(0.6), Inches(0.25), Inches(12), Inches(0.5), "Layout 重构：从 600 行 if-else 到 6 个子类的多态分发", sz=Pt(24), color=BLUE, bold=True)

code_box(s, Inches(0.3), Inches(0.85), Inches(6.0), [
    "# 当前: if-else 分发 (~600行)",
    "def _reshape(...):",
    "    if self.use_compress:        # DS V4",
    "        ...",
    "    elif self.use_sparse:        # DS V3.2/GLM5.1",
    "        if C8 and self.is_A5:   # A5 sparse C8",
    "            ...",
    "        elif C8:                # A3 sparse C8",
    "            ...",
    "        else:                   # no C8",
    "            ...",
    "    elif self.use_hybrid_blocks: # Qwen3.5",
    "        ...",
    "    elif self.cache_only_layers:",
    "        ...",
    "    else:                        # GQA / standard MLA",
    "        ...",
], sz=Pt(7))

code_box(s, Inches(6.6), Inches(0.85), Inches(6.1), [
    "# 目标: 多态分发 (~135行 + 625行Layout类)",
    "class KVCacheLayout(ABC):",
    "    def get_kv_cache_shape() -> list[tuple]",
    "    def split_sizes(total,spec) -> list[int]",
    "    def reshape(raws,spec) -> list[Tensor]",
    "",
    "class StandardKVLayout       # GQA → 2 tensor",
    "class MLANopeRopeLayout     # DS V3.1 → 2 tensor",
    "class SparseMLALayout       # DS V3.2/GLM5.1 → 3~4t",
    "class CompressedMLALayout   # DS V4 → as_strided views",
    "class HybridLayout          # Qwen3.5 → shared buffer",
    "class CacheOnlyLayout       # hidden state → 1 tensor",
    "",
    "# SparseMLALayout(precision='fp8_c8', device='A5')",
    "# → 构造参数消化精度+设备差异，对 model_runner 透明",
], sz=Pt(7))

ref_data = [
    ["指标", "旧代码", "新代码"],
    ["model_runner行数", "~600 行 if-else", "~135行调度 + 625行Layout类"],
    ["新增模型方式", "找所有分支点 + 改3~5处", "创建/选择Layout子类 → 实现3方法"],
    ["单元测试", "难以独立测试各分支", "每个Layout子类独立可测"],
    ["分支状态", "feature/layout-refactor-phase3", "12 files, +2080/-3357 lines"],
]
make_table(s, Inches(0.3), Inches(4.2), Inches(12.4), Inches(1.5), ref_data)

card(s, Inches(0.3), Inches(5.9), Inches(12.4), Inches(0.9),
     "重构后 model_runner 的核心逻辑 (伪代码)", [
         "for layer_name, spec in kv_cache_specs.items():",
         "    layout = spec.get_layout()     # 多态: 根据 spec 类型+字段自动选择子类",
         "    shapes = layout.get_kv_cache_shape(N, B, H, D)  # 返回 list[shape], 不是单个 shape",
         "    sizes  = layout.split_sizes(total_bytes, spec)  # 返回 list[int], 每个 sub-tensor 的字节数",
         "    raw_tensors = [torch.zeros(s, dtype=int8) for s in sizes]  # 分配",
         "    kv_caches[name] = layout.reshape(raw_tensors, spec)         # 重塑",
         "→ model_runner 不再有 if-else。Layout 子类内部消化所有差异。",
     ], tc=GREEN, accent=GREEN)

# ═══════════════════════════════════════════
# SLIDE 29 — 已对齐部分详解
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_PRIMARY); slide_num(s, 29, 30)
tb(s, Inches(0.6), Inches(0.25), Inches(12), Inches(0.5), "社区已对齐机制 + 仍需覆盖的 Gap", sz=Pt(24), color=BLUE, bold=True)

# Left: aligned
aligned_data = [
    ["已对齐机制", "说明", "对齐程度"],
    ["KVCacheSpec 基类", "page_size_bytes / max_memory_usage 基础计算", "100%"],
    ["KVCacheConfig", "num_blocks + kv_cache_tensors + kv_cache_groups", "100%"],
    ["KVCacheTensor.shared_by", "多层共享物理 buffer 机制", "100%"],
    ["BlockPool", "block 分配/回收/哈希管理", "100%"],
    ["SingleTypeKVCacheManager", "Full/Sliding/Chunked/Mamba 管理器", "100%"],
    ["KVCacheCoordinator", "Unitary/Hybrid 协调器 (仅DS V4有patch)", "95%"],
    ["Prefix Caching", "hash 计算 + block 管理 + 跨请求复用", "100%"],
    ["get_kv_cache_groups()", "决策级联 (仅DS V4有patch覆盖)", "90%"],
]
make_table(s, Inches(0.3), Inches(0.8), Inches(6.5), Inches(2.5), aligned_data, header_color='C6000B')

# Right: gaps
gap_data = [
    ["仍需覆盖的 Gap", "差距描述", "优先级"],
    ["Spec 层: 万能类问题", "所有MLA变体挤进 AscendMLAAttentionSpec\n缺少按类型多态的 Spec 子类", "P0"],
    ["Backend 接口: 1/N 信息", "get_kv_cache_shape() 只返回参考shape\n需要扩展为 list[tuple] 或 Layout 对象", "P0"],
    ["Allocate: if-else 决策树", "185行, 深度3, 每次加模型改3~5处\n需要 Layout.split_sizes() 替代", "P0"],
    ["Reshape: 精度+设备分支", "350行, A5/A3+C8/非C8+sparse/compress\n需要 Layout.reshape() 替代", "P0"],
    ["DS V4: 分桶逻辑 patch", "patch_kv_cache_utils 覆盖分组函数\n需确认上游 DS V4 支持进度", "P1"],
    ["算子 stride 支持", "长期目标: 减少 host 侧预拆分需求\n需要与算子侧持续对齐推进", "P2"],
]
make_table(s, Inches(7.1), Inches(0.8), Inches(6.0), Inches(2.3), gap_data, header_color='C8002E')

card(s, Inches(0.3), Inches(3.5), Inches(6.5), Inches(3.2),
     "为什么已对齐部分可以做到 100%", [
         "这些组件属于「规划层」——在 tensor 实际分配之前",
         "负责: 分组决策、block 管理、哈希计算、配置生成",
         "不涉及: NPU 特有的 stride 限制、多 tensor 拆分",
         "→ 与硬件无关，Ascend 可以直接使用上游代码",
         "",
         "而「执行层」(Allocate + Reshape) 差异巨大:",
         "上游: 拿到 KVCacheConfig → torch.zeros → view → done",
         "Ascend: 拿到 KVCacheConfig → 检查 5 个标志 → 拆 tensor",
         "  → 算各自 shape/dtype → view → 组装 tuple",
         "→ 这正是 KVCacheLayout 要封装的层次",
     ], tc=GREEN, accent=GREEN)

card(s, Inches(7.1), Inches(3.5), Inches(5.8), Inches(3.2),
     "DS V4 patch 详细说明", [
         "patch_kv_cache_utils.py 覆盖了两个函数:",
         "1. resolve_kv_cache_block_sizes():",
         "   处理 CP (Context Parallel) 下 block_size 缩放",
         "   覆盖上游 PR #40860 对多 group + CP 的限制",
         "2. group_and_unify_kv_cache_specs():",
         "   DS V4 的 Compressor+Indexer+SWA 三层结构",
         "   需要按 page_size 分桶 + layer_tuple 对齐",
         "",
         "patch_kv_cache_coordinator.py:",
         "   AscendHybridKVCacheCoordinator",
         "   覆盖 find_longest_cache_hit() 逻辑",
         "",
         "评估: 上游已支持 DS V4 → 这些 patch",
         "  未来可通过 Layout 重构合并或删除",
     ], tc=AMBER, accent=AMBER)

# ═══════════════════════════════════════════
# SLIDE 30 — 总结 (重构版)
# ═══════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_PRIMARY); slide_num(s, 30, 30)
tb(s, Inches(0.6), Inches(0.15), Inches(12), Inches(0.35),
   "总结：三大差异 → 三大对策 → 三步路线 → 一套架构",
   sz=Pt(26), color=BLUE, bold=True)
tb(s, Inches(0.6), Inches(0.48), Inches(12), Inches(0.22),
   "从诊断到方案到执行：每一层差异都有对应策略，每一步都有明确交付物",
   sz=Pt(10), color=TEXT_SEC)

# ── Row 1: Problem (left) vs Solution (right) ──
card(s, Inches(0.3), Inches(0.82), Inches(6.4), Inches(2.25),
     "三大差异 (诊断)", [
         "① 算子约束: 上游 GPU kernel 接受 blob 内部解析 → Ascend NPU 不支持 stride",
         "   → host 侧被迫预拆分 tensor → 1 tensor 变 N tensors",
         "",
         "② 类型擦除: 上游 isinstance(spec) 编码全部语义 → Ascend 万能类 if-else 恢复",
         "   → 5 种 MLA 变体共用一个 AscendMLAAttentionSpec → 下游 15+ 分支",
         "",
         "③ 接口不足: 上游 get_kv_cache_shape() 返回完整信息 → Ascend 同名方法返回 1/N",
         "   → model_runner 绕过 Backend 直连 attention layer 掏参数 → ~600 行逻辑",
     ], tc=RED, accent=RED)

card(s, Inches(7.0), Inches(0.82), Inches(6.0), Inches(2.25),
     "三大对策 (方案)", [
         "① KVCacheLayout 多态: 6 个子类封装 NPU 多 tensor 拆分逻辑",
         "   layout.split_sizes() + layout.reshape() → model_runner 零分支遍历",
         "",
         "② Spec 子类拆分: 4~5 个子类替代万能类, 恢复类型即语义",
         "   isinstance(spec, AscendSparseC8MLASpec) → 不检查字段就知道全部信息",
         "   spec.get_layout() → 自动选择对应 Layout → Spec-Layout 互补",
         "",
         "③ Backend 接口扩展: get_kv_cache_shape() → list[tuple] 或 get_layout()",
         "   → model_runner 不再需要绕过后端直取 attention layer 属性",
     ], tc=GREEN, accent=GREEN)

# ── Row 2: Final architecture ──
card(s, Inches(0.3), Inches(3.30), Inches(12.4), Inches(1.2),
     "对齐后的最终架构: 端到端 4 层 → model_runner ~60 行, 0 分支", [
         "layer.get_kv_cache_spec()  →  KVCacheSpec (纯声明)  →  isinstance() 编码全部语义",
         "spec.get_layout()           →  KVCacheLayout 子类   →  多态选择, 零 if-else, 封装 NPU 拆分逻辑",
         "layout.get_kv_cache_shape() →  list[tuple]           →  每个 sub-tensor 的完整 shape, 不再 1/N",
         "layout.split_sizes()        →  list[int]             →  model_runner 无分支分配各 sub-tensor",
         "layout.reshape(raws, spec)  →  list[Tensor]          →  dtype 正确, 直接传入 NPU 算子",
     ], tc=GREEN, accent=GREEN)

# ── Row 3: 3-phase timeline ──
phases = [
    ("Phase 1\nQ3", GREEN,
     "内部重构\nSpec 4~5 子类\nLayout 6 子类\nmodel_runner\n~135 行"),
    ("Phase 2\nQ4", BLUE,
     "上游 RFC\nLayout 抽象提议\nBackend 接口扩展\nSpec 字段\n通用化"),
    ("Phase 3\nQ1+", PURPLE,
     "社区合入\n+ 清理\n删除 patch\n删除 Ascend\n万能类"),
]
for i, (title, clr, desc) in enumerate(phases):
    x = Inches(0.3) + Inches(i*4.3)
    # Phase badge
    badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(4.75), Inches(1.5), Pt(36))
    r2 = max(0, min(255, int(clr[0]*0.1))); g2 = max(0, min(255, int(clr[1]*0.1))); b2 = max(0, min(255, int(clr[2]*0.1)))
    badge.fill.solid(); badge.fill.fore_color.rgb = RGBColor(r2,g2,b2); badge.line.fill.background()
    badge.adjustments[0] = 0.3
    p = badge.text_frame.paragraphs[0]; p.text = title; p.font.size = Pt(11); p.font.color.rgb = clr; p.font.bold = True
    p.font.name = 'Microsoft YaHei'; p.alignment = PP_ALIGN.CENTER
    # Description
    tb(s, x + Inches(1.7), Inches(4.75), Inches(2.3), Pt(36), desc, sz=Pt(8), color=TEXT_SEC)
    # Arrow between phases
    if i < 2:
        tb(s, x + Inches(4.0), Inches(4.83), Inches(0.3), Pt(20), "→", sz=Pt(16), color=RED, bold=True, align=PP_ALIGN.CENTER)

# Bottom metrics
tb(s, Inches(0.6), Inches(5.55), Inches(12), Inches(0.2),
   "目标: model_runner 回归上游的简洁 — 不检查 layer_type, 不判断 C8, 不区分 A5/A3, 不知道精度",
   sz=Pt(10), color=RED, bold=True)

stat(s, Inches(1.5), Inches(5.85), "~600→60", "model_runner 行数", GREEN)
stat(s, Inches(4.7), Inches(5.85), "15+→0", "if-else 分支", GREEN)
stat(s, Inches(7.9), Inches(5.85), "6→1", "通用路径", GREEN)
stat(s, Inches(11.0), Inches(5.85), "独立可测", "Layout 子类", GREEN)

# Divider bar at bottom
bar2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.0), prs.slide_width, Pt(5))
bar2.fill.solid(); bar2.fill.fore_color.rgb = PRIMARY_RED; bar2.line.fill.background()

# ═══════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════
output = r"c:\Users\c50058674\code\vllm-ascend\docs\kv_cache_presentation.pptx"
prs.save(output)
print(f"Done! {len(prs.slides)} slides saved to {output}")
