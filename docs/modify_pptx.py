#!/usr/bin/env python3
"""Modify the user's PPTX with 3 improvements:
1. Slide 22: Change root cause #3 + "两个源头"→"三个源头"
2. Slide 14: Add upstream reference columns to comparison table
3. Add transition slide after Slide 8 (community section)
"""

import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
import copy

# ── Color constants ──
PRIMARY_RED = RGBColor(198, 0, 11)
RED_TINT3 = RGBColor(248, 210, 212)
GREEN = RGBColor(88, 178, 48)
AMBER = RGBColor(237, 108, 0)
MAGENTA = RGBColor(198, 0, 84)
BLUE = RGBColor(198, 0, 11)  # Same as PRIMARY_RED in this scheme
BLACK = RGBColor(0, 0, 0)
DARK_GRAY = RGBColor(89, 87, 87)
TEXT_SEC = DARK_GRAY
OFF_WHITE = RGBColor(245, 245, 247)
WHITE = RGBColor(255, 255, 255)
BORDER = RGBColor(218, 218, 222)
BG_CARD = OFF_WHITE
TEXT_PRIMARY = RGBColor(30, 30, 30)
BG_PRIMARY = WHITE
RED = RGBColor(200, 16, 46)
PURPLE = MAGENTA

prs = Presentation('vllm-ascend-KV cache管理优化方案.pptx')

# ═══════════════════════════════════════════════════════
# IMPROVEMENT 1: Slide 22 — Root cause #3 + "三个源头"
# ═══════════════════════════════════════════════════════
slide22 = prs.slides[21]

for shape in slide22.shapes:
    if shape.has_text_frame:
        tf = shape.text_frame
        full_text = tf.text

        # Fix card #3 title
        if '3. Sparse' in full_text:
            for p in tf.paragraphs:
                if p.text.startswith('3. Sparse'):
                    for run in p.runs:
                        run.text = run.text.replace('3. Sparse 多 Tensor', '3. Spec 万能类: 类型擦除')
                        run.font.color.rgb = MAGENTA
                    break
                elif ('aclnnSparseFlashAttention' in p.text or
                      '3~4独立tensor输入' in p.text or
                      'host侧预拆分' in p.text):
                    for run in p.runs:
                        run.text = '5种MLA变体→1个AscendMLAAttentionSpec。isinstance()不编码语义，下游用if-else恢复已丢失的类型信息。'
                    break

        # Fix callout title: "两个根本源头" → "三个根本源头"
        if '两个根本源头' in full_text:
            for p in tf.paragraphs:
                if '两个根本源头' in p.text:
                    for run in p.runs:
                        if '两个' in run.text:
                            run.text = run.text.replace('两个', '三个')
                    break

        # Fix callout body
        if '硬件约束: NPU算子不支持stride' in full_text and '接口缺陷:' in full_text:
            for p in tf.paragraphs:
                if p.text and ('硬件约束:' in p.text or '硬件约束：' in p.text):
                    for run in p.runs:
                        if '接口缺陷' in run.text:
                            run.text = (
                                '硬件约束: NPU算子不支持stride访问 → host侧必须预拆分tensor → 分配和重塑逻辑复杂化 (根因1,2,4,5)。'
                                'Spec 设计: 类型不编码语义 → isinstance()无区分力 → 下游被迫用if-else恢复丢失的类型信息 (根因3)。'
                                '接口缺陷: get_kv_cache_shape() 单返回值无法表达多tensor布局 → 信息泄漏到 model_runner (根因6)。'
                            )
                    break

# Also fix the layout refactor card: add Spec subclass mention
for shape in slide22.shapes:
    if shape.has_text_frame:
        tf = shape.text_frame
        if 'KVCacheLayout 抽象基类' in tf.text and 'StandardKV' in tf.text:
            # This is the Layout refactor card - update it
            for p in tf.paragraphs:
                if '6 个子类' in p.text:
                    for run in p.runs:
                        if '6 个子类' in run.text:
                            run.text = run.text.replace(
                                '6 个子类: StandardKV, MLANopeRope, SparseMLA,',
                                'Spec 4~5 子类 + Layout 6 个子类:\n  Spec: AscendStandardMLA/SparseMLA/SparseC8MLA/DSV4MLA\n  Layout: StandardKV, MLANopeRope, SparseMLA,'
                            )
                    break

print('✓ Improvement 1: Slide 22 updated.')

# ═══════════════════════════════════════════════════════
# IMPROVEMENT 2: Slide 14 — Add upstream reference columns
# ═══════════════════════════════════════════════════════
slide14 = prs.slides[13]

# Find the table
for shape in slide14.shapes:
    if shape.has_table:
        tbl = shape.table
        rows, cols = len(tbl.rows), len(tbl.columns)

        # Check if this is the model comparison table
        header_text = tbl.cell(0, 0).text
        if '模型类型' in header_text or '模型' in header_text:
            # Read current data
            old_data = []
            for r in range(rows):
                row_data = []
                for c in range(cols):
                    row_data.append(tbl.cell(r, c).text)
                old_data.append(row_data)

            # Build new table data with upstream columns
            new_data = [
                ["模型类型", "上游 Spec\n(基线)", "Ascend Spec\n(额外字段)",
                 "上游 Backend\n(基线)", "Ascend Backend",
                 "上游 Reshape", "Ascend Reshape\n(额外复杂度)", "差距 Δ"],
                ["GQA\n(Qwen3 MoE)",
                 "FullAttentionSpec\n7字段,零NPU",
                 "复用上游 ✓\n无额外字段",
                 "FlashAttention\n→ (2,N,B,H,D)",
                 "AscendAttention\n→ (2,N,B,H,D)",
                 "view(N,B,H,D)\n1行",
                 ".view()\n同上游",
                 "★\nK/V物理拆分"],
                ["标准MLA\n(DS V3.1)",
                 "MLAAttentionSpec\nhead_size=576",
                 "AscendMLAAttention\n+scale_dim/dtype",
                 "FlashMLA\n→ (N,B,576)",
                 "AscendMLA\n→ (N,B,1,576)",
                 "view(N,B,576)\n1行",
                 "view+查layer\nhead_size丢失512+64",
                 "★★\nhead_size=576\n丢失拆分信息"],
                ["SparseMLA\n(DS V3.2)",
                 "MLAAttentionSpec\n+cache_dtype_str",
                 "AscendMLAAttention\n+sparse_head_dim/C8",
                 "FlashMLASparse\n→ (N,B,576)",
                 "AscendSFA\n→ (N,B,1,576)",
                 "view(N,B,576)\n1行",
                 "unpack+view\nC8×A5/A3多分支",
                 "★★★\nindexer+C8量化\nA5/A3硬件分支"],
                ["CompressMLA\n(DS V4)",
                 "MLAAttentionSpec\n+compress_ratio",
                 "AscendMLAAttention\n+compress_ratio/scale",
                 "FlashMLASparse\n→ (N,B,656)",
                 "AscendDSA\n→ (N,B,1,head)",
                 "view(N,B,656)\n1行",
                 "as_strided overlay\n3 views",
                 "★★★★\n时间压缩+fp8\nas_strided承载语义"],
                ["Hybrid\n(Qwen3.5)",
                 "FullAttn+MambaSpec\n2种原生类型",
                 "复用上游 ✓\n+page_size_padded",
                 "FlashAttn+GDN\n→ 两种shape",
                 "AscendAttn+GDN\n→ (2,N,B,4,256)",
                 "view+as_strided_\n零拷贝重排",
                 "slice+strip+padding\n~100行hybrid专用",
                 "★★★★★\nattn+mamba\nlayout互斥"],
            ]

            # Rebuild table
            new_rows, new_cols = len(new_data), len(new_data[0])

            # Get table position
            tbl_left = shape.left
            tbl_top = shape.top
            tbl_width = shape.width
            tbl_height = shape.height

            # Remove old table
            sp = shape._element
            sp.getparent().remove(sp)

            # Create new table
            new_shape = slide14.shapes.add_table(new_rows, new_cols,
                                                  tbl_left, tbl_top,
                                                  Inches(12.7), Inches(3.4))
            new_tbl = new_shape.table

            # Style the table
            for r in range(new_rows):
                for c in range(new_cols):
                    cell = new_tbl.cell(r, c)
                    p = cell.text_frame.paragraphs[0]
                    p.text = new_data[r][c]
                    p.font.size = Pt(6.5)
                    p.font.color.rgb = TEXT_PRIMARY
                    p.font.name = 'Consolas' if r > 0 else 'Microsoft YaHei'

                    if r == 0:
                        p.font.bold = True
                        p.font.color.rgb = WHITE
                        p.font.size = Pt(7)

                    # Cell fill
                    tcPr = cell._tc.get_or_add_tcPr()
                    fill = cell._tc.makeelement(qn('a:solidFill'), {})
                    if r == 0:
                        clr_val = 'C6000B'
                    else:
                        clr_val = 'F5F5F7'
                    clr = cell._tc.makeelement(qn('a:srgbClr'), {'val': clr_val})
                    fill.append(clr)
                    tcPr.insert(0, fill)

            # Update card titles below the table to reflect new messaging
            for shape2 in slide14.shapes:
                if shape2.has_text_frame:
                    if shape2.top > Inches(4.0):  # Cards below the table
                        tf2 = shape2.text_frame
                        if '上游基线' in tf2.text and '走同一条' in tf2.text:
                            for p2 in tf2.paragraphs:
                                if '整个 pipeline' in p2.text:
                                    for run in p2.runs:
                                        if '不感知模型差异' in run.text:
                                            run.text = run.text.replace(
                                                '不感知模型差异',
                                                '不感知模型差异 → 零 NPU 概念'
                                            )

            print('✓ Improvement 2: Slide 14 table rebuilt with upstream columns.')
            break

# ═══════════════════════════════════════════════════════
# IMPROVEMENT 3: Add transition slide after Slide 8
# ═══════════════════════════════════════════════════════
# Find available slide layouts
layouts = prs.slide_layouts
# Use the first blank-ish layout (usually index 6 is blank)
# Try layouts that exist
target_layout = None
for layout in layouts:
    if layout.name and 'blank' in layout.name.lower():
        target_layout = layout
        break
if target_layout is None:
    # Try layout 6 or the last layout
    if len(layouts) > 6:
        target_layout = layouts[6]
    else:
        target_layout = layouts[-1]

# Create the transition slide
new_slide = prs.slides.add_slide(target_layout)

# Set white background
bg = new_slide.background
bg.fill.solid()
bg.fill.fore_color.rgb = WHITE

# Build the slide content
def add_textbox(slide, l, t, w, h, text, sz=Pt(14), color=TEXT_PRIMARY,
                bold=False, align=PP_ALIGN.LEFT, font='Microsoft YaHei'):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = sz
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return tf

def add_card(slide, l, t, w, h, title, lines, tc=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_CARD
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(1)
    shape.adjustments[0] = 0.04

    # Accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, Pt(3), h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = tc
    bar.line.fill.background()

    add_textbox(slide, l+Pt(12), t+Pt(8), w-Pt(24), Pt(22),
                title, sz=Pt(14), color=tc, bold=True)

    box = slide.shapes.add_textbox(l+Pt(12), t+Pt(34), w-Pt(24), h-Pt(42))
    tf = box.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = ""
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(9.5)
        p.font.color.rgb = TEXT_SEC
        p.font.name = 'Microsoft YaHei'
        p.line_spacing = Pt(17)

# Title
add_textbox(new_slide, Inches(0.6), Inches(2.2), Inches(12), Inches(0.7),
            '社区 vLLM KV Cache 管线 — 三个核心设计原则',
            sz=Pt(28), color=PRIMARY_RED, bold=True, align=PP_ALIGN.CENTER)

add_textbox(new_slide, Inches(0.6), Inches(2.9), Inches(12), Inches(0.4),
            '记住这三条原则，就能理解为什么 Ascend 的差异不是"实现细节不同"，而是"设计原则层面的偏离"',
            sz=Pt(11), color=TEXT_SEC, align=PP_ALIGN.CENTER)

# Three principle cards
principles = [
    ("原则 1: Backend 说了算", GREEN, [
        "KV Cache shape 的唯一决定者是 Backend",
        "get_kv_cache_shape() 返回完整可用 shape",
        "每一维的含义、tensor 的数量、K/V 的存储方式——都由 Backend 封装",
        "→ model_runner 不关心、也不需要知道内部怎么拆分",
    ]),
    ("原则 2: 单 Tensor 原则", AMBER, [
        "所有模型类型都产生 1 个 tensor/层",
        "（或 1 个 buffer + as_strided views）",
        "MLA: kv_lora + k_rope 在同一 tensor 内，kernel 自己 offset",
        "Mamba: 1 buffer → N 个 as_strided 零拷贝视图",
        "→ Host 侧只管分配，不关心内部数据组织",
    ]),
    ("原则 3: model_runner 零感知", RED, [
        "model_runner 不检查 layer_type",
        "不判断 use_sparse / use_compress / C8 / A5 / A3",
        "拿到 Spec → 拿到 shape → 分配 → reshape → 结束",
        "→ 所有模型特有逻辑封装在 Spec 子类 + Backend 内部",
        "→ ~60 行代码，0 分支，新增模型不需要改 model_runner",
    ]),
]

for i, (title, clr, lines) in enumerate(principles):
    x = Inches(0.3) + Inches(i * 4.3)
    add_card(new_slide, x, Inches(3.5), Inches(4.0), Inches(3.5),
             title, lines, tc=clr)

# Bottom line
add_textbox(new_slide, Inches(0.6), Inches(7.1), Inches(12), Inches(0.3),
            '➜ 接下来我们将逐层对比 Ascend 的实现，看看在每一层上，这三条原则是如何被打破的',
            sz=Pt(11), color=RED, bold=True, align=PP_ALIGN.CENTER)

# Move the new slide to position 8 (after slide 8, index 7)
# Get the XML element for slide ordering
sldIdLst = prs.part.element.find(qn('p:sldIdLst'))
sldId_elements = list(sldIdLst)

# The new slide is currently the last one
new_sldId = sldId_elements[-1]

# Remove from end
sldIdLst.remove(new_sldId)

# Insert after the 8th slide (index 7, which is now the 8th element after removal)
# Wait - after removing from end, the indices shift
# The original slide 8 was at index 7. It's still at index 7 since we removed from end.
target = sldId_elements[7]  # This is still the original slide 8
# Insert after it
target_index = list(sldIdLst).index(target)
sldIdLst.insert(target_index + 1, new_sldId)

print('✓ Improvement 3: Transition slide added after Slide 8.')

# ═══════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════
output_path = 'vllm-ascend-KV cache管理优化方案.pptx'
prs.save(output_path)
print(f'\n✓ PPTX saved to: {output_path}')
print(f'Total slides: {len(prs.slides)}')
